#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Hindi Script to PowerPoint Converter
=====================================

यह program Hindi TTS script को PowerPoint presentation में convert करता है

Author: Your Name
License: MIT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os
import sys

# ============================================================================
# CONFIGURATION - यहाँ settings change कर सकते हैं
# ============================================================================

class Config:
    """Configuration settings for PPT generation"""
    
    # File Settings
    INPUT_FILE = 'input_script.txt'
    OUTPUT_FILE = 'output_presentation.pptx'
    
    # Slide Settings
    SLIDE_WIDTH = Inches(10)    # 16:9 ratio
    SLIDE_HEIGHT = Inches(5.625)
    
    # Colors (RGB format) - Modern Professional Palette
    TITLE_COLOR = RGBColor(26, 35, 126)      # Deep Blue
    SUBTITLE_COLOR = RGBColor(67, 160, 71)   # Green
    HEADING_COLOR = RGBColor(13, 71, 161)    # Blue
    BODY_COLOR = RGBColor(33, 33, 33)        # Dark Gray
    BG_COLOR = RGBColor(255, 255, 255)       # White
    ACCENT_COLOR = RGBColor(255, 152, 0)     # Orange
    GRADIENT_START = RGBColor(232, 245, 233) # Light Green
    GRADIENT_END = RGBColor(227, 242, 253)   # Light Blue
    
    # Font Settings
    TITLE_FONT = 'Calibri'
    BODY_FONT = 'Calibri'
    
    TITLE_FONT_SIZE = Pt(48)
    SUBTITLE_FONT_SIZE = Pt(30)
    HEADING_FONT_SIZE = Pt(36)
    BODY_FONT_SIZE = Pt(20)
    BULLET_FONT_SIZE = Pt(18)
    
    # Layout Settings
    MAX_BULLETS_PER_SLIDE = 8  # More content per slide
    MAX_CHARS_PER_BULLET = 200  # Longer bullets allowed

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def print_status(message, symbol="ℹ️"):
    """Print formatted status message"""
    print(f"{symbol} {message}")

def create_title_slide(prs, title_text, subtitle_text):
    """
    Title slide बनाता है
    
    Args:
        prs: Presentation object
        title_text: Main title
        subtitle_text: Subtitle text
    """
    print_status("Creating title slide...", "📄")
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add colored bar at top
    top_bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(0.4)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = Config.HEADING_COLOR
    top_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), 
        Inches(9), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Config.TITLE_FONT_SIZE
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = Config.TITLE_COLOR
    title_frame.paragraphs[0].font.name = Config.TITLE_FONT
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.8), 
        Inches(8), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle_text
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_frame.paragraphs[0].font.size = Config.SUBTITLE_FONT_SIZE
    subtitle_frame.paragraphs[0].font.color.rgb = Config.SUBTITLE_COLOR
    subtitle_frame.paragraphs[0].font.name = Config.BODY_FONT
    
    # Add bottom accent line
    bottom_line = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(3), Inches(4.5),
        Inches(4), Inches(0.05)
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = Config.ACCENT_COLOR
    bottom_line.line.fill.background()
    
    return slide

def create_section_slide(prs, section_title):
    """
    Section title slide बनाता है
    
    Args:
        prs: Presentation object
        section_title: Section का title
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Section title (centered)
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2), 
        Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = section_title
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Config.HEADING_FONT_SIZE
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = Config.HEADING_COLOR
    title_frame.paragraphs[0].font.name = Config.TITLE_FONT
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return slide

def create_content_slide(prs, title, content_lines):
    """
    Content slide बनाता है with title and bullet points
    
    Args:
        prs: Presentation object
        title: Slide का title
        content_lines: List of content lines
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), 
        Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = Config.HEADING_COLOR
    title_frame.paragraphs[0].font.name = Config.TITLE_FONT
    
    # Content
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), 
        Inches(9), Inches(4)
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = line
        p.font.size = Config.BODY_FONT_SIZE
        p.font.color.rgb = Config.BODY_COLOR
        p.font.name = Config.BODY_FONT
        p.space_after = Pt(10)
        p.level = 0
    
    return slide

def create_bullet_slide(prs, title, bullets):
    """
    Bullet points वाली slide बनाता है with modern design
    
    Args:
        prs: Presentation object
        title: Slide का title
        bullets: List of bullet points
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add colored header bar
    header_bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(0.9)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = Config.HEADING_COLOR
    header_bar.line.fill.background()
    
    # Title on header bar
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15), 
        Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(34)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White
    title_frame.paragraphs[0].font.name = Config.TITLE_FONT
    
    # Bullets
    content_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.2), 
        Inches(8.7), Inches(4.2)
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = "● " + bullet  # Use filled circle
        p.font.size = Config.BULLET_FONT_SIZE
        p.font.color.rgb = Config.BODY_COLOR
        p.font.name = Config.BODY_FONT
        p.space_after = Pt(12)
        p.level = 0
    
    return slide

def create_end_slide(prs):
    """
    अंतिम slide बनाता है with modern design
    
    Args:
        prs: Presentation object
    """
    print_status("Creating end slide...", "📄")
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background accent
    accent_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(5.625)
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = Config.GRADIENT_START
    accent_shape.line.fill.background()
    
    # Thank You text
    end_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2), 
        Inches(7), Inches(1.5)
    )
    end_frame = end_box.text_frame
    end_frame.text = "Thank You!"
    end_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    end_frame.paragraphs[0].font.size = Pt(60)
    end_frame.paragraphs[0].font.bold = True
    end_frame.paragraphs[0].font.color.rgb = Config.HEADING_COLOR
    end_frame.paragraphs[0].font.name = Config.TITLE_FONT
    end_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(2), Inches(3.5), 
        Inches(6), Inches(0.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Questions?"
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = Config.SUBTITLE_COLOR
    subtitle_frame.paragraphs[0].font.name = Config.BODY_FONT
    subtitle_frame.paragraphs[0].font.italic = True
    
    return slide

def split_into_chunks(lines, max_per_slide=6):
    """
    Content को multiple slides में split करता है
    
    Args:
        lines: List of content lines
        max_per_slide: Maximum lines per slide
    
    Returns:
        List of chunks (each chunk is a list of lines)
    """
    chunks = []
    current_chunk = []
    
    for line in lines:
        current_chunk.append(line)
        if len(current_chunk) >= max_per_slide:
            chunks.append(current_chunk)
            current_chunk = []
    
    if current_chunk:
        chunks.append(current_chunk)
    
    return chunks

# ============================================================================
# MAIN PROCESSING FUNCTION
# ============================================================================

def process_script(input_file, output_file):
    """
    Main function to process script and create PPT
    
    Args:
        input_file: Path to input script file
        output_file: Path to output PPT file
    """
    
    # Check if input file exists
    if not os.path.exists(input_file):
        print_status(f"Error: Input file '{input_file}' not found!", "❌")
        print_status(f"Tip: अपनी script को '{input_file}' नाम से save करें", "💡")
        return False
    
    print_status("Reading input file...", "📖")
    
    # Read script
    try:
        with open(input_file, 'r', encoding='utf-8') as f:
            content = f.read()
    except Exception as e:
        print_status(f"Error reading file: {e}", "❌")
        return False
    
    lines = content.split('\n')
    
    print_status("Processing script...", "✍️")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Config.SLIDE_WIDTH
    prs.slide_height = Config.SLIDE_HEIGHT
    
    # Extract title and subtitle
    title_text = "Hindi Script Presentation"
    subtitle_text = ""
    
    for line in lines[:5]:
        line = line.strip()
        if line and 'Text-to-Speech' not in line:
            if title_text == "Hindi Script Presentation":
                title_text = line.split(':')[0].strip()
            elif not subtitle_text:
                if line.startswith('(') and line.endswith(')'):
                    subtitle_text = line[1:-1]
                elif not line.startswith('==='):
                    subtitle_text = line
                break
    
    # Create title slide
    create_title_slide(prs, title_text, subtitle_text)
    
    # Process sections
    skip_header = True
    current_section = None
    section_content = []
    section_count = 0
    slide_count = 1  # Title slide already created
    
    for line in lines:
        line = line.strip()
        
        # Skip header
        if skip_header:
            if 'Text-to-Speech' in line:
                skip_header = False
            continue
        
        # Skip empty lines and markers
        if not line or line == '[समाप्त]':
            continue
        
        # Section header
        if line.startswith('===') and 'भाग' in line:
            # Save previous section
            if current_section and section_content:
                section_count += 1
                
                # Create section title slide
                create_section_slide(prs, current_section)
                slide_count += 1
                
                # Split content into multiple slides if needed
                chunks = split_into_chunks(section_content, Config.MAX_BULLETS_PER_SLIDE)
                
                for i, chunk in enumerate(chunks):
                    slide_title = current_section
                    if len(chunks) > 1:
                        slide_title += f" ({i+1}/{len(chunks)})"
                    
                    create_content_slide(prs, slide_title, chunk)
                    slide_count += 1
            
            # Start new section
            current_section = line.replace('===', '').strip()
            section_content = []
            continue
        
        # Add content to current section
        if current_section:
            section_content.append(line)
    
    # Process last section
    if current_section and section_content:
        section_count += 1
        create_section_slide(prs, current_section)
        slide_count += 1
        
        chunks = split_into_chunks(section_content, Config.MAX_BULLETS_PER_SLIDE)
        for i, chunk in enumerate(chunks):
            slide_title = current_section
            if len(chunks) > 1:
                slide_title += f" ({i+1}/{len(chunks)})"
            create_content_slide(prs, slide_title, chunk)
            slide_count += 1
    
    # Create end slide
    create_end_slide(prs)
    slide_count += 1
    
    # Save presentation
    print_status("Saving presentation...", "💾")
    try:
        prs.save(output_file)
    except Exception as e:
        print_status(f"Error saving file: {e}", "❌")
        return False
    
    print()
    print_status("SUCCESS!", "✅")
    print_status(f"Presentation created: {output_file}", "📄")
    print_status(f"Total slides: {slide_count}", "📊")
    print_status(f"Total sections: {section_count}", "📑")
    print()
    print_status("अब आप PowerPoint file को open कर सकते हैं!", "🎉")
    
    return True

# ============================================================================
# MAIN EXECUTION
# ============================================================================

def main():
    """Main function"""
    
    print()
    print("═" * 50)
    print("  Hindi Script to PowerPoint Converter")
    print("═" * 50)
    print()
    
    # Process script
    success = process_script(Config.INPUT_FILE, Config.OUTPUT_FILE)
    
    if not success:
        print()
        print_status("Troubleshooting tips:", "💡")
        print("  1. Check if input file exists")
        print("  2. Make sure file is UTF-8 encoded")
        print("  3. Check file permissions")
        print()
        sys.exit(1)

if __name__ == "__main__":
    main()
