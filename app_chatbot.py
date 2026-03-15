import streamlit as st

import os
import re
import time
import subprocess
import tempfile
import json as json_mod
from datetime import datetime

# Auto-install node_modules on Streamlit Cloud (npm not available locally but nodejs is)
_node_pptx_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "node_pptx")
if not os.path.exists(os.path.join(_node_pptx_dir, "node_modules")):
    try:
        subprocess.run(["npm", "install"], cwd=_node_pptx_dir, capture_output=True, timeout=120)
    except Exception:
        pass
from document_upload_component import document_upload_component

# Imports
try:
    from ppt_to_images import ppt_to_images
    PPT_TO_IMAGES_AVAILABLE = True
except:
    PPT_TO_IMAGES_AVAILABLE = False

from content_generator import generate_content_from_topic
from ai_ppt_generator import generate_beautiful_ppt, create_chart_image
from multi_ai_generator import MultiAIGenerator, get_last_ai_source
from web_search import search_google

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

google_api_key = os.getenv("GOOGLE_API_KEY")
google_cse_id = os.getenv("GOOGLE_CSE_ID")

# Page Config
st.set_page_config(
    page_title="FREE PPT Maker - AI Presentation Generator",
    page_icon="📊",
    layout="centered",
    initial_sidebar_state="collapsed"
)

# Modern ChatGPT-style CSS with Light/Dark Mode Support
st.markdown("""
<style>
    /* ═══════════════════════════════════════════════════════════════════
       🎨 CSS VARIABLES FOR THEMING - Easy Light/Dark Mode Support
       ═══════════════════════════════════════════════════════════════════ */

    :root {
        /* Light Mode Colors (Default) */
        --bg-primary: #ffffff;
        --bg-secondary: #f4f6f9;
        --bg-tertiary: #e8ecf4;
        --text-primary: #1a1a2e;
        --text-secondary: #4a4a6a;
        --text-muted: #6b7280;
        --accent-primary: #6366f1;
        --accent-secondary: #8b5cf6;
        --accent-gradient: linear-gradient(135deg, #6366f1 0%, #8b5cf6 50%, #06b6d4 100%);
        --border-color: #d1d5db;
        --border-focus: #6366f1;
        --shadow-sm: 0 2px 8px rgba(0, 0, 0, 0.08);
        --shadow-md: 0 4px 16px rgba(0, 0, 0, 0.12);
        --shadow-lg: 0 8px 32px rgba(0, 0, 0, 0.16);
        --btn-text: #ffffff;
        --input-bg: #ffffff;
        --card-bg: #ffffff;
        --success-bg: #d1fae5;
        --success-text: #065f46;
    }

    /* Dark Mode Colors */
    @media (prefers-color-scheme: dark) {
        :root {
            --bg-primary: #0f172a;
            --bg-secondary: #1e293b;
            --bg-tertiary: #334155;
            --text-primary: #f1f5f9;
            --text-secondary: #cbd5e1;
            --text-muted: #94a3b8;
            --accent-primary: #818cf8;
            --accent-secondary: #a78bfa;
            --accent-gradient: linear-gradient(135deg, #818cf8 0%, #a78bfa 50%, #22d3ee 100%);
            --border-color: #475569;
            --border-focus: #818cf8;
            --shadow-sm: 0 2px 8px rgba(0, 0, 0, 0.3);
            --shadow-md: 0 4px 16px rgba(0, 0, 0, 0.4);
            --shadow-lg: 0 8px 32px rgba(0, 0, 0, 0.5);
            --btn-text: #ffffff;
            --input-bg: #1e293b;
            --card-bg: #1e293b;
            --success-bg: #064e3b;
            --success-text: #6ee7b7;
        }
    }

    /* Hide Streamlit defaults */
    #MainMenu, footer, [data-testid="stHeader"] {display: none !important;}

    /* ═══════════════════════════════════════════════════════════════════
       🎯 BASE STYLES
       ═══════════════════════════════════════════════════════════════════ */

    body, .stApp, .main, .block-container {
        background: var(--bg-secondary) !important;
        color: var(--text-primary) !important;
    }

    .block-container {
        padding-top: 90px !important;
        padding-left: 2rem !important;
        padding-right: 2rem !important;
        padding-bottom: 220px !important;
        max-width: 1000px !important;
        margin: 0 auto !important;
    }

    /* ═══════════════════════════════════════════════════════════════════
       📱 MOBILE RESPONSIVE STYLES
       ═══════════════════════════════════════════════════════════════════ */

    @media (max-width: 768px) {
        .block-container {
            padding-top: 80px !important;
            padding-left: 0.5rem !important;
            padding-right: 0.5rem !important;
            padding-bottom: 200px !important;
        }

        .header-container {
            padding: 12px 16px;
            border-radius: 0 0 14px 14px;
        }

        .header-title {
            font-size: 22px !important;
            letter-spacing: 0.5px;
        }

        .header-subtitle {
            font-size: 10px;
        }

        [data-testid="stChatMessageContainer"] {
            padding: 12px 10px !important;
            border-radius: 12px;
        }

        .stChatMessage {
            padding: 0.8rem !important;
            font-size: 15px !important;
            border-radius: 12px !important;
        }

        .stChatInput {
            width: 95% !important;
            bottom: 70px !important;
            padding: 10px 12px !important;
            border-radius: 14px !important;
        }

        .stChatInput textarea,
        [data-testid="stChatInput"] textarea {
            min-height: 50px !important;
            font-size: 14px !important;
        }

        .stButton > button,
        .stDownloadButton > button {
            font-size: 14px !important;
            padding: 10px 16px !important;
        }

        .footer-container {
            padding: 10px 12px !important;
            font-size: 11px !important;
            bottom: 0 !important;
        }

        .icon-btn-container button {
            width: 40px !important;
            height: 40px !important;
            min-width: 40px !important;
            max-width: 40px !important;
            font-size: 16px !important;
        }
    }

    @media (max-width: 480px) {
        .header-title {
            font-size: 18px !important;
        }

        .block-container {
            padding-top: 70px !important;
        }

        .stChatInput {
            bottom: 60px !important;
        }
    }

    /* ═══════════════════════════════════════════════════════════════════
       🦶 FOOTER STYLES
       ═══════════════════════════════════════════════════════════════════ */

    .footer-container {
        position: fixed;
        bottom: 0;
        left: 0;
        right: 0;
        background: var(--accent-gradient);
        color: #ffffff;
        text-align: center;
        padding: 8px 16px;
        font-size: 13px;
        z-index: 9998;
        box-shadow: 0 -2px 10px rgba(0,0,0,0.15);
        display: flex;
        justify-content: center;
        align-items: center;
        gap: 20px;
    }

    .footer-container span {
        opacity: 0.9;
    }

    .footer-stat {
        background: rgba(255,255,255,0.2);
        padding: 2px 10px;
        border-radius: 12px;
        font-weight: 600;
        font-size: 12px;
    }

    /* ═══════════════════════════════════════════════════════════════════
       🌈 HEADER SECTION
       ═══════════════════════════════════════════════════════════════════ */

    .header-container {
        background: var(--accent-gradient);
        padding: 16px 24px;
        border-radius: 0 0 20px 20px;
        margin-bottom: 28px;
        box-shadow: var(--shadow-lg);
        text-align: center;
        width: 100%;
        position: fixed;
        top: 0;
        left: 0;
        right: 0;
        z-index: 9999;
    }

    .header-title {
        font-size: 32px;
        font-weight: 900;
        color: #ffffff !important;
        margin: 0;
        line-height: 1.1;
        text-shadow: 0 2px 12px rgba(0,0,0,0.4);
        letter-spacing: 1px;
    }

    .header-subtitle {
        font-size: 12px;
        color: rgba(255,255,255,0.85);
        margin-top: 4px;
        font-weight: 400;
    }

    /* ═══════════════════════════════════════════════════════════════════
       💬 CHAT AREA STYLES
       ═══════════════════════════════════════════════════════════════════ */

    [data-testid="stChatMessageContainer"] {
        margin-bottom: 2rem !important;
        min-height: 350px;
        max-width: 900px;
        margin-left: auto;
        margin-right: auto;
        background: var(--card-bg) !important;
        border-radius: 20px;
        box-shadow: var(--shadow-md);
        border: 1px solid var(--border-color);
        padding: 24px 32px;
    }

    .stChatMessage {
        margin-bottom: 1.2rem !important;
        padding: 1.2rem !important;
        border-radius: 16px !important;
        max-width: 800px;
        margin-left: auto;
        margin-right: auto;
        font-size: 17px !important;
        color: var(--text-primary) !important;
        background: var(--bg-tertiary) !important;
        box-shadow: var(--shadow-sm);
        border: 1px solid var(--border-color);
    }

    .stChatMessage.user {
        background: var(--accent-primary) !important;
        color: #ffffff !important;
        border-top-right-radius: 6px !important;
        border-bottom-left-radius: 22px !important;
        border: none;
    }

    .stChatMessage.assistant {
        background: var(--card-bg) !important;
        color: var(--text-primary) !important;
        border-top-left-radius: 6px !important;
        border-bottom-right-radius: 22px !important;
    }

    /* ═══════════════════════════════════════════════════════════════════
       📝 CHAT INPUT STYLES
       ═══════════════════════════════════════════════════════════════════ */

    .stChatInput {
        position: fixed !important;
        bottom: 30px !important;
        left: 50% !important;
        transform: translateX(-50%) !important;
        width: 90% !important;
        max-width: 900px !important;
        z-index: 100 !important;
        background: var(--card-bg) !important;
        border: 2px solid var(--accent-primary) !important;
        border-radius: 18px !important;
        padding: 16px 20px !important;
        box-shadow: var(--shadow-lg) !important;
    }

    .stChatInput textarea,
    [data-testid="stChatInput"] textarea {
        min-height: 80px !important;
        font-size: 16px !important;
        padding: 12px !important;
        background: var(--input-bg) !important;
        color: var(--text-primary) !important;
        border-radius: 12px !important;
        border: 2px solid var(--border-color) !important;
        transition: border-color 0.2s, box-shadow 0.2s;
    }

    .stChatInput textarea:focus,
    [data-testid="stChatInput"] textarea:focus {
        border-color: var(--accent-primary) !important;
        box-shadow: 0 0 0 3px rgba(99, 102, 241, 0.2) !important;
        outline: none !important;
    }

    .stChatInput textarea::placeholder {
        color: var(--text-muted) !important;
        opacity: 1 !important;
    }

    /* ═══════════════════════════════════════════════════════════════════
       🔘 BUTTON STYLES - HIGH VISIBILITY
       ═══════════════════════════════════════════════════════════════════ */

    /* Primary Buttons (Download, Submit, etc.) */
    .stButton > button,
    .stDownloadButton > button,
    button[data-testid="stBaseButton-primary"] {
        background: var(--accent-gradient) !important;
        color: #ffffff !important;
        font-weight: 700 !important;
        font-size: 16px !important;
        border-radius: 12px !important;
        box-shadow: var(--shadow-md) !important;
        padding: 12px 24px !important;
        border: none !important;
        transition: all 0.25s ease !important;
        text-shadow: 0 1px 2px rgba(0,0,0,0.2);
    }

    .stButton > button:hover,
    .stDownloadButton > button:hover,
    button[data-testid="stBaseButton-primary"]:hover {
        transform: translateY(-2px) !important;
        box-shadow: var(--shadow-lg) !important;
        filter: brightness(1.1);
    }

    /* Secondary Buttons (Skip, Cancel, etc.) */
    button[data-testid="stBaseButton-secondary"] {
        background: var(--bg-tertiary) !important;
        color: var(--text-primary) !important;
        font-weight: 600 !important;
        font-size: 15px !important;
        border: 2px solid var(--accent-primary) !important;
        border-radius: 12px !important;
        padding: 10px 20px !important;
        transition: all 0.25s ease !important;
    }

    button[data-testid="stBaseButton-secondary"]:hover {
        background: var(--accent-primary) !important;
        color: #ffffff !important;
        transform: translateY(-2px) !important;
    }

    /* Icon Buttons ONLY - Round style for single emoji buttons */
    .icon-btn-container button {
        background: var(--card-bg) !important;
        border: 2px solid var(--accent-primary) !important;
        border-radius: 50% !important;
        width: 48px !important;
        height: 48px !important;
        min-width: 48px !important;
        max-width: 48px !important;
        font-size: 20px !important;
        padding: 0 !important;
        box-shadow: var(--shadow-md) !important;
        cursor: pointer;
        transition: all 0.25s ease !important;
        color: var(--accent-primary) !important;
    }

    .icon-btn-container button:hover {
        background: var(--accent-primary) !important;
        color: #ffffff !important;
        transform: scale(1.1) !important;
        box-shadow: var(--shadow-lg) !important;
    }

    /* ═══════════════════════════════════════════════════════════════════
       📋 INPUT FIELDS - CLEAR & VISIBLE
       ═══════════════════════════════════════════════════════════════════ */

    /* Text Inputs */
    .stTextInput input,
    .stTextInput input[type="text"],
    input[data-testid="stTextInput"] {
        background: var(--input-bg) !important;
        color: var(--text-primary) !important;
        border: 2px solid var(--border-color) !important;
        border-radius: 12px !important;
        font-size: 16px !important;
        font-weight: 500 !important;
        padding: 14px 16px !important;
        transition: all 0.2s ease !important;
        box-shadow: var(--shadow-sm) !important;
    }

    .stTextInput input:focus {
        border-color: var(--accent-primary) !important;
        box-shadow: 0 0 0 4px rgba(99, 102, 241, 0.25) !important;
        outline: none !important;
    }

    .stTextInput input::placeholder {
        color: var(--text-muted) !important;
        opacity: 1 !important;
        font-weight: 400 !important;
    }

    /* Labels */
    .stTextInput label,
    label[data-testid="stWidgetLabel"] {
        color: var(--text-secondary) !important;
        font-weight: 600 !important;
        font-size: 14px !important;
        margin-bottom: 8px !important;
    }

    /* Select Box / Dropdown */
    .stSelectbox > div > div,
    [data-testid="stSelectbox"] {
        background: var(--input-bg) !important;
        color: var(--text-primary) !important;
        border: 2px solid var(--border-color) !important;
        border-radius: 12px !important;
    }

    .stSelectbox [data-baseweb="select"] > div {
        background: var(--input-bg) !important;
        border-color: var(--border-color) !important;
    }

    /* ═══════════════════════════════════════════════════════════════════
       📎 FILE UPLOAD & MISC ELEMENTS
       ═══════════════════════════════════════════════════════════════════ */

    .file-attached {
        background: var(--success-bg) !important;
        color: var(--success-text) !important;
        padding: 10px 16px;
        border-radius: 10px;
        font-weight: 600;
        display: inline-flex;
        align-items: center;
        gap: 8px;
        border: 1px solid var(--success-text);
    }

    /* Success/Error Messages */
    .stSuccess {
        background: var(--success-bg) !important;
        color: var(--success-text) !important;
        border-radius: 12px !important;
        border: 1px solid var(--success-text) !important;
    }

    /* Markdown text */
    .stMarkdown, .stMarkdown p {
        color: var(--text-primary) !important;
    }

    /* Spinner/Loading */
    .stSpinner > div {
        border-top-color: var(--accent-primary) !important;
    }

    /* Hide floating action bar */
    .floating-actions-bar { display: none !important; }

    /* ═══════════════════════════════════════════════════════════════════
       📱 ACTION ROW ABOVE CHAT INPUT
       ═══════════════════════════════════════════════════════════════════ */

    .chatbox-action-row {
        display: flex;
        flex-direction: row;
        gap: 12px;
        justify-content: flex-start;
        align-items: center;
        margin: 0 auto 12px auto;
        max-width: 900px;
        padding: 0 8px;
    }
</style>
""", unsafe_allow_html=True)

# ═══════════════════════════════════════════════════════════════════════════
# 🌈 COLORFUL HEADER SECTION
# ═══════════════════════════════════════════════════════════════════════════

st.markdown('<div class="main-container">', unsafe_allow_html=True)
st.markdown('<div class="header-container"><h1 class="header-title">FREE PPT Generator</h1><div class="header-subtitle">AI-Powered Presentation Maker | No Login Required</div></div>', unsafe_allow_html=True)
st.markdown('<div class="content-container">', unsafe_allow_html=True)


# Initialize session state
if 'messages' not in st.session_state:
    st.session_state.messages = []
if 'stage' not in st.session_state:
    st.session_state.stage = 'idle'  # idle, ask_name, ask_designation, awaiting_topic, awaiting_theme, confirming, generating, done
    # Start with welcome message first, ask name later if needed
if 'presenter_name' not in st.session_state:
    st.session_state.presenter_name = None
if 'presenter_designation' not in st.session_state:
    st.session_state.presenter_designation = None
if 'topic' not in st.session_state:
    st.session_state.topic = None
if 'suggested_title' not in st.session_state:
    st.session_state.suggested_title = None
if 'ppt_path' not in st.session_state:
    st.session_state.ppt_path = None
if 'theme' not in st.session_state:
    st.session_state.theme = 'modern'
if 'bullets_per_slide' not in st.session_state:
    st.session_state.bullets_per_slide = 4  # Default: 4 bullets per slide
if 'file_content' not in st.session_state:
    st.session_state.file_content = None
if 'file_name' not in st.session_state:
    st.session_state.file_name = None
if 'ai_source' not in st.session_state:
    st.session_state.ai_source = None
if 'num_slides' not in st.session_state:
    st.session_state.num_slides = 6
if 'parsed_slides' not in st.session_state:
    st.session_state.parsed_slides = None
if 'file_names' not in st.session_state:
    st.session_state.file_names = []
if 'file_contents' not in st.session_state:
    st.session_state.file_contents = []
if 'chart_data' not in st.session_state:
    st.session_state.chart_data = None  # Stores {'df': DataFrame, 'filename': str}
if 'chart_settings' not in st.session_state:
    st.session_state.chart_settings = None  # Stores {'chart_type': str, 'columns': list}

# Helper functions
def add_message(role, content):
    st.session_state.messages.append({"role": role, "content": content})

def is_greeting(text):
    # Global greetings - English, Hindi, Spanish, French, German, Arabic, Chinese, Japanese, etc.
    greetings = [
        # English
        'hi', 'hello', 'hey', 'hii', 'hiii', 'hiiii', 'helo', 'hllo', 'hlw', 'hw',
        'good morning', 'good afternoon', 'good evening', 'good night', 'gm', 'gn',
        'howdy', 'yo', 'sup', 'whats up', "what's up", 'wassup', 'wazzup',
        'greetings', 'salutations', 'heya', 'hiya', 'hai',
        # Hindi / Hinglish
        'namaste', 'namaskar', 'pranam', 'pranaam', 'namashkar', 'jai shri krishna',
        'jai hind', 'radhe radhe', 'ram ram', 'jai shri ram', 'jai mata di',
        'kya hal', 'kaise ho', 'kaise hai', 'kaisa hai', 'kya haal hai', 'kya chal raha',
        'aur batao', 'sab theek', 'theek hai', 'haan bhai', 'bhai', 'bro',
        'sat sri akal', 'sasriakal', 'kem cho', 'vanakkam', 'nomoshkar',
        # Spanish
        'hola', 'buenos dias', 'buenas tardes', 'buenas noches', 'que tal',
        # French
        'bonjour', 'bonsoir', 'salut', 'coucou',
        # German
        'hallo', 'guten tag', 'guten morgen', 'guten abend',
        # Italian
        'ciao', 'buongiorno', 'buonasera',
        # Portuguese
        'ola', 'bom dia', 'boa tarde', 'boa noite',
        # Arabic
        'marhaba', 'ahlan', 'salam', 'assalamualaikum', 'as-salamu alaykum', 'salaam',
        # Chinese
        'ni hao', 'nihao',
        # Japanese
        'konnichiwa', 'ohayo', 'konbanwa',
        # Korean
        'annyeong', 'annyeonghaseyo',
        # Russian
        'privet', 'zdravstvuyte',
        # Other
        'aloha', 'shalom', 'sawadee', 'jambo', 'habari'
    ]
    text_lower = text.lower().strip()
    # Remove punctuation for matching
    text_clean = re.sub(r'[^\w\s]', '', text_lower)
    for g in greetings:
        g_clean = re.sub(r'[^\w\s]', '', g)
        if text_clean == g_clean or text_clean.startswith(g_clean + ' '):
            return True
    return False

def is_yes(text):
    yes_words = ['yes', 'ya', 'yaa', 'haan', 'han', 'ok', 'okay', 'sure', 'theek', 'thik', 'sahi', 'correct', 'proceed', 'go ahead', 'bana do', 'banao']
    return text.lower().strip() in yes_words or text.lower().strip().startswith('yes')

def is_no(text):
    no_words = ['no', 'nahi', 'naa', 'na', 'change', 'modify', 'edit', 'different']
    return text.lower().strip() in no_words

def is_valid_topic(text):
    """Check if the input is a valid/meaningful topic for PPT generation"""
    text = text.strip()

    # ISSUE 10: Minimum 10 characters for short inputs
    if len(text) < 10:
        return False, "Please provide more details (at least 10 characters). Example: 'AI in Healthcare' or 'Digital India'."

    # ISSUE 7: Long/structured content is always valid (user pasted content)
    # If text is substantial (> 50 chars), has multiple lines, or has bullet points - accept it
    if len(text) > 50:
        return True, ""

    # Multi-line content is likely structured/pasted content
    if '\n' in text and len(text.split('\n')) >= 2:
        return True, ""

    # Content with bullet markers is structured
    if any(marker in text for marker in ['- ', '• ', '* ', '1. ', '2. ']):
        return True, ""

    # Check for minimum vowels (gibberish often lacks vowels)
    vowels = set('aeiouAEIOU')
    vowel_count = sum(1 for c in text if c in vowels)

    # If mostly consonants with very few vowels, likely gibberish
    if len(text) > 5 and vowel_count < len(text) * 0.15:
        return False, "Yeh samajh nahi aaya. Please ek clear topic likhen jaise 'Artificial Intelligence' ya 'Climate Change'."

    # Check for repeated characters (like "aaaaaaa" or "jjjjj")
    if len(text) > 4:
        for i in range(len(text) - 3):
            if text[i] == text[i+1] == text[i+2] == text[i+3]:
                return False, "Please ek meaningful topic enter karein, random characters nahi."

    # Check if it's just numbers
    if text.replace(' ', '').isdigit():
        return False, "Please topic ka naam likhen, sirf numbers nahi."

    # If text is long enough and has some vowels, probably valid
    if len(text) >= 10 and vowel_count >= 2:
        return True, ""

    return False, "Please ek clear topic batayein (minimum 10 characters). Example: 'Digital India' ya 'AI in Healthcare'."

def fix_pptx_backgrounds(pptx_path):
    """Post-process PPTX: replace slide.background <p:bgPr> with full-slide rectangle shapes.
    Fixes blank display in PowerPoint/Google Slides where layout bg1=white overrides bgPr.
    """
    import zipfile as _zipfile, shutil as _shutil, re as _re
    import os as _os

    # Slide dimensions for LAYOUT_WIDE (13.33" x 7.5" in EMU)
    SLIDE_W = 12192000
    SLIDE_H = 6858000

    try:
        with _zipfile.ZipFile(pptx_path, 'r') as zin:
            info_map = {info.filename: info for info in zin.infolist()}
            file_contents = {info.filename: zin.read(info.filename) for info in zin.infolist()}

        modified = False
        slide_pattern = _re.compile(r'^ppt/slides/slide\d+\.xml$')
        bg_pattern = _re.compile(
            r'<p:bg><p:bgPr><a:solidFill><a:srgbClr val="([0-9A-Fa-f]{6})"/>'
            r'</a:solidFill></p:bgPr></p:bg>'
        )

        for fname in file_contents:
            if not slide_pattern.match(fname):
                continue
            xml = file_contents[fname].decode('utf-8')

            # Skip if already processed
            if 'name="BG_RECT"' in xml:
                continue

            match = bg_pattern.search(xml)
            if not match:
                continue

            bg_color = match.group(1).upper()

            # Build full-slide background rectangle as first shape in spTree
            bg_rect = (
                f'<p:sp><p:nvSpPr><p:cNvPr id="999" name="BG_RECT"/>'
                f'<p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
                f'<p:spPr><a:xfrm><a:off x="0" y="0"/>'
                f'<a:ext cx="{SLIDE_W}" cy="{SLIDE_H}"/></a:xfrm>'
                f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
                f'<a:solidFill><a:srgbClr val="{bg_color}"/></a:solidFill>'
                f'<a:ln><a:noFill/></a:ln></p:spPr></p:sp>'
            )

            # Insert after </p:grpSpPr> (before other shapes = bottom of z-order)
            insert_after = '</p:grpSpPr>'
            pos = xml.find(insert_after)
            if pos == -1:
                continue
            pos += len(insert_after)
            xml = xml[:pos] + bg_rect + xml[pos:]
            file_contents[fname] = xml.encode('utf-8')
            modified = True

        if not modified:
            return

        # Rewrite PPTX preserving original compress_type per file
        tmp_path = pptx_path + '.bgfix.tmp'
        with _zipfile.ZipFile(tmp_path, 'w') as zout:
            for arcname, content in file_contents.items():
                orig_info = info_map.get(arcname)
                compress = orig_info.compress_type if orig_info else _zipfile.ZIP_DEFLATED
                zout.writestr(arcname, content, compress_type=compress)

        _shutil.move(tmp_path, pptx_path)
        print(f"[FIX_BG] Background rectangles added to PPTX slides.")

    except Exception as e:
        print(f"[FIX_BG] Post-processing failed (non-critical): {e}")
        try:
            if _os.path.exists(pptx_path + '.bgfix.tmp'):
                _os.unlink(pptx_path + '.bgfix.tmp')
        except:
            pass


def run_pptxgenjs(js_code, output_path):
    """Execute PptxGenJS code via Node.js subprocess. Returns (success, path_or_error)."""
    project_dir = os.path.dirname(os.path.abspath(__file__))
    wrapper_path = os.path.join(project_dir, "node_pptx", "pptx_wrapper.js")
    node_cwd = os.path.join(project_dir, "node_pptx")

    # Strip any pptx/PptxGenJS redeclarations the AI may have added
    import re as _re2
    js_code = _re2.sub(r'^\s*(const|let|var)\s+pptx\s*=\s*new\s+PptxGenJS[^\n]*\n?', '', js_code, flags=_re2.MULTILINE)
    js_code = _re2.sub(r'^\s*(const|let|var)\s+pptx\s*=\s*require[^\n]*\n?', '', js_code, flags=_re2.MULTILINE)
    js_code = _re2.sub(r'^\s*pptx\.layout\s*=[^\n]*\n?', '', js_code, flags=_re2.MULTILINE)

    # Write AI code to temp file
    fd, temp_js_path = tempfile.mkstemp(suffix='.js', prefix='pptx_slides_')
    try:
        with os.fdopen(fd, 'w') as f:
            f.write(js_code)

        result = subprocess.run(
            ['node', wrapper_path, output_path, temp_js_path],
            capture_output=True, text=True, timeout=45,
            cwd=node_cwd
        )

        if result.returncode == 0:
            try:
                response = json_mod.loads(result.stdout.strip())
                if response.get('success'):
                    # Post-process: add background rectangles for viewer compatibility
                    fix_pptx_backgrounds(response['path'])
                    return True, response['path']
            except:
                pass

        # Parse error
        err_msg = result.stderr or result.stdout or 'Node.js execution failed'
        try:
            err_data = json_mod.loads(err_msg.strip())
            return False, err_data.get('error', err_msg)
        except:
            return False, err_msg
    except subprocess.TimeoutExpired:
        return False, 'Node.js execution timed out (45s)'
    except Exception as e:
        return False, str(e)
    finally:
        try:
            os.unlink(temp_js_path)
        except:
            pass


def generate_ppt(content, topic, theme):
    """Generate PPT — tries PptxGenJS first, falls back to python-pptx."""
    project_dir = os.path.dirname(os.path.abspath(__file__))
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_folder = os.path.join(project_dir, "output", f"output_{timestamp}")
    os.makedirs(output_folder, exist_ok=True)

    # Generate meaningful filename from topic
    if topic:
        words = re.sub(r'[^\w\s]', '', str(topic)).split()
        meaningful_words = [w for w in words[:5] if len(w) > 2][:4]
        topic_slug = '_'.join(meaningful_words) if meaningful_words else 'presentation'
    else:
        topic_slug = 'presentation'
    topic_slug = re.sub(r'[^\w]', '_', topic_slug)[:40]
    topic_slug = re.sub(r'_+', '_', topic_slug).strip('_')
    date_str = datetime.now().strftime("%Y%m%d")
    ppt_filename = f"{topic_slug}_{date_str}.pptx"
    ppt_path = os.path.join(output_folder, ppt_filename)

    # ── PRIMARY: PptxGenJS via Node.js ──
    js_code = st.session_state.get('pptxgenjs_code')
    print(f"[GENERATE_PPT] js_code present: {bool(js_code)}, length: {len(js_code) if js_code else 0}")
    if js_code:
        success, result = run_pptxgenjs(js_code, ppt_path)
        print(f"[GENERATE_PPT] run_pptxgenjs result: success={success}, result={str(result)[:200]}")
        if success:
            print(f"[PPTXGENJS] Successfully generated: {ppt_path}")
            return True, ppt_path
        else:
            print(f"[PPTXGENJS] Failed: {result}, falling back to python-pptx")

    # ── FALLBACK: python-pptx ──
    # Guard: if content is empty, nothing to render → fail gracefully
    has_content = isinstance(content, list) and any(
        slide.get('bullets') or slide.get('main_title') or slide.get('title', '').strip()
        for slide in content if isinstance(slide, dict)
    )
    if not has_content:
        print("[GENERATE_PPT] No slide content to render — AI generation failed.")
        return False, "AI generation failed. Please try again."
    # Inject chart slide if Excel/CSV data is available
    if isinstance(content, list) and st.session_state.get('chart_data') and PANDAS_AVAILABLE:
        chart_data = st.session_state.chart_data
        chart_settings = st.session_state.get('chart_settings', {})
        chart_type = chart_settings.get('chart_type', 'bar')
        chart_title = chart_settings.get('chart_title', f"Data from {chart_data['filename']}")
        chart_img = create_chart_image(chart_data['df'], chart_type=chart_type, title=chart_title)
        if chart_img:
            chart_slide = {"slide_number": 2, "title": chart_title, "bullets": [], "chart_image_path": chart_img}
            content.insert(1, chart_slide)

    if isinstance(content, list) and all(isinstance(slide, dict) for slide in content):
        success = generate_beautiful_ppt(content, ppt_path, color_scheme=theme, use_ai=False, original_topic=topic, min_slides=6, max_slides=6, generate_ai_images=True)
    else:
        success = generate_beautiful_ppt(content, ppt_path, color_scheme=theme, use_ai=False, original_topic=topic, min_slides=6, max_slides=6, generate_ai_images=True)
    return success, ppt_path

# ═══════════════════════════════════════════════════════════════════════════════
# 🔧 SLIDE EDIT DETECTION & REGENERATION FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════
def detect_slide_edit_request(user_text):
    """
    Detect if user wants to edit a specific slide.
    Returns: (slide_number, edit_instruction) or (None, None)
    """
    user_text_lower = user_text.lower()

    # Patterns to detect slide number
    patterns = [
        r'slide\s*(\d+)',           # "slide 3", "slide3"
        r'स्लाइड\s*(\d+)',          # Hindi: "स्लाइड 3"
        r'(\d+)\s*(st|nd|rd|th)\s*slide',  # "3rd slide"
        r'#\s*(\d+)',               # "#3"
    ]

    slide_num = None
    for pattern in patterns:
        match = re.search(pattern, user_text_lower)
        if match:
            slide_num = int(match.group(1))
            break

    if slide_num:
        return slide_num, user_text
    return None, None

def regenerate_single_slide(slide_num, instruction, all_slides, topic, language):
    """
    Regenerate a single slide based on user instruction.
    """
    if slide_num < 1 or slide_num > len(all_slides):
        return None, "Invalid slide number"

    current_slide = all_slides[slide_num - 1]

    try:
        generator = MultiAIGenerator()

        prompt = f"""You need to modify Slide {slide_num} of a presentation on "{topic}".

Current slide content:
Title: {current_slide.get('title', '')}
Bullets: {current_slide.get('bullets', [])}

User's instruction: {instruction}

STRICT RULES:
- Slide title: max 35 characters
- Each bullet: exactly 1 COMPLETE sentence, max 15 words
- Never truncate or use "..." — write short but complete sentences
- 4 bullet points total

Generate the updated slide in this exact format:
Slide {slide_num}: [Short Title max 35 chars]
- Complete sentence, max 15 words.
- Complete sentence, max 15 words.
- Complete sentence, max 15 words.
- Complete sentence, max 15 words.

Language: {language}
Only output the slide content, nothing else."""

        content_dict = generator.generate_ppt_content(
            topic=f"Modify slide {slide_num}",
            min_slides=1,
            max_slides=1,
            custom_instructions=prompt,
            bullets_per_slide=st.session_state.get('bullets_per_slide', 4),
            bullet_word_limit=25
        )

        ai_output = content_dict.get("output", "")
        if not ai_output:
            return None, "AI returned empty response"

        # Parse the single slide
        lines = [l.strip() for l in ai_output.split('\n') if l.strip()]
        new_slide = {"slide_number": slide_num, "bullets": []}

        for line in lines:
            m = re.match(r"^\*{0,2}Slide\s*\d+\s*[:\-–]\s*(.+?)\*{0,2}$", line, re.IGNORECASE)
            if m:
                new_slide["title"] = m.group(1).strip()
            elif line.startswith('- ') or line.startswith('• ') or line.startswith('* '):
                new_slide["bullets"].append(line[2:].strip())
            elif re.match(r'^\d+[\.\)]\s+', line):
                new_slide["bullets"].append(re.sub(r'^\d+[\.\)]\s+', '', line).strip())

        if new_slide.get("title") or new_slide.get("bullets"):
            # Preserve original slide properties not being changed
            if not new_slide.get("title"):
                new_slide["title"] = current_slide.get("title", f"Slide {slide_num}")
            return new_slide, None
        else:
            return None, "Could not parse AI response"

    except Exception as e:
        return None, str(e)

# Additional CSS - Using CSS variables (defined above)

# Welcome bar at top (helps visibility on mobile - user suggested)

# Remove Streamlit native header and caption; only show custom HTML header

# Language selection (bullets are now flexible 4-6 based on content)
language = st.selectbox("Language", ["English", "Hindi"], key="ppt_language", index=0, label_visibility="collapsed")
st.session_state.language = language
# Bullets per slide is now flexible (4-6) based on content needs - AI decides automatically
st.session_state.bullets_per_slide = 5  # Default/average for compatibility

# Display chat messages
for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])

# Auto-scroll to bottom when new messages appear
if st.session_state.messages:
    st.markdown("""
    <script>
        var chatContainer = window.parent.document.querySelector('[data-testid="stChatMessageContainer"]');
        if (chatContainer) {
            chatContainer.scrollTop = chatContainer.scrollHeight;
        }
        // Fallback: scroll main container
        var main = window.parent.document.querySelector('.main');
        if (main) {
            main.scrollTop = main.scrollHeight;
        }
    </script>
    """, unsafe_allow_html=True)

# Chat input (fixed at bottom by Streamlit)
user_input = st.chat_input("Send topic, paste content (Hindi/English)...", key="main_chat_input")

# Show download button if PPT is ready
if st.session_state.stage == 'done' and st.session_state.ppt_path:
    with st.chat_message("assistant"):
        if st.session_state.topic and st.session_state.topic.lower() != 'none':
            st.success(f"Your presentation on **{st.session_state.topic}** is ready!")
        col1, col2 = st.columns([3, 1])
        with col1:
            download_filename = os.path.basename(st.session_state.ppt_path) if st.session_state.ppt_path else "presentation.pptx"
            with open(st.session_state.ppt_path, "rb") as f:
                st.download_button("⬇️ Download PPT", f.read(), file_name=download_filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True, type="primary")
        with col2:
            if st.button("New", use_container_width=True):
                st.session_state.messages = []
                st.session_state.stage = 'idle'
                st.session_state.ppt_path = None
                st.session_state.topic = None
                st.session_state.file_content = None
                st.session_state.parsed_slides = None
                st.session_state.file_names = []
                st.session_state.file_contents = []
                st.rerun()

# Chat input

# Name and designation collection flow
if st.session_state.stage == 'ask_name':
    with st.chat_message("assistant"):
        st.markdown("**Enter your name (for the presentation): (optional)**")
        col1, col2 = st.columns([3,1])
        with col1:
            name = st.text_input("Your name", key="presenter_name_input", label_visibility="visible", placeholder="Type your name...", help="This will appear on the PPT (optional)")
        with col2:
            skip_btn = st.button("Skip", key="skip_name", help="Skip name", use_container_width=True)
            if skip_btn:
                st.session_state.presenter_name = None
                st.session_state.stage = 'ask_designation'
                st.rerun()
        if name:
            st.session_state.presenter_name = name
            st.session_state.stage = 'ask_designation'
            st.rerun()
elif st.session_state.stage == 'ask_designation':
    with st.chat_message("assistant"):
        st.markdown("**Enter your designation (for the presentation): (optional)**")
        col1, col2 = st.columns([3,1])
        with col1:
            designation = st.text_input("Your designation", key="presenter_designation_input", label_visibility="visible", placeholder="Type your designation...", help="This will appear on the PPT (optional)")
        with col2:
            skip_btn = st.button("Skip", key="skip_designation", help="Skip designation", use_container_width=True)
            if skip_btn:
                st.session_state.presenter_designation = None
                st.session_state.stage = 'idle'
                # Show personalized welcome
                name = st.session_state.presenter_name
                if name:
                    welcome = f"Welcome {name}! Please tell me a topic for your presentation."
                else:
                    welcome = "Welcome! Please tell me a topic for your presentation."
                add_message("assistant", welcome)
                st.rerun()
        if designation:
            st.session_state.presenter_designation = designation
            st.session_state.stage = 'idle'
            name = st.session_state.presenter_name
            if name:
                welcome = f"Welcome {name}! Please tell me a topic for your presentation."
            else:
                welcome = "Welcome! Please tell me a topic for your presentation."
            add_message("assistant", welcome)
            st.rerun()

# ============ MODERN CHATBOX UI ============
# Initialize show_uploader state
if 'show_uploader' not in st.session_state:
    st.session_state.show_uploader = False

# Store all uploaded file names and contents in session state as lists
if 'file_names' not in st.session_state:
    st.session_state.file_names = []
if 'file_contents' not in st.session_state:
    st.session_state.file_contents = []

# File attached indicator (show above chatbox)
if st.session_state.get('file_names'):
    chips = []
    for idx, name in enumerate(st.session_state['file_names']):
        chip_col, btn_col = st.columns([8,1], gap="small")
        with chip_col:
            st.markdown(f'<div class="file-attached" style="display: flex; align-items: center; background: #e0f7fa; color: #00796b; border-radius: 16px; padding: 2px 10px 2px 6px; font-size: 14px; font-weight: 500; box-shadow: 0 1px 3px rgba(0,0,0,0.04); border: 1px solid #b2dfdb; max-width: 220px; overflow: hidden; white-space: nowrap; margin-bottom: 0;">'
                        f'<span style="font-size:16px; margin-right:4px;">📎</span>'
                        f'<span style="overflow: hidden; text-overflow: ellipsis; max-width: 120px; display: inline-block;">{name}</span></div>', unsafe_allow_html=True)
        with btn_col:
            if st.button("✕", key=f"clear_file_{idx}", help=f"Remove {name}", use_container_width=True):
                del st.session_state.file_names[idx]
                del st.session_state.file_contents[idx]
                if not st.session_state.file_names:
                    st.session_state.file_content = None
                    st.session_state.file_name = None
                    st.session_state.uploaded_preview = None
                st.rerun()
# Floating left action buttons (Upload, Refresh) - Top left above chat

# Floating left action buttons (Upload, Refresh) - Fixed top left above chat

# --- Insert action icons inside chat input (left of arrow) ---

# --- True inline action icons inside chat input (left of textarea, inside box) ---
 # --- Desktop-friendly chat input with action icons row above input ---
# (CSS already defined in main style block above)

# Icon buttons row (Upload +, Refresh 🔄)
st.markdown('<div class="icon-btn-container" style="display: flex; gap: 12px; margin-bottom: 16px;">', unsafe_allow_html=True)
icon_col1, icon_col2, icon_spacer = st.columns([1, 1, 10])
with icon_col1:
    upload_clicked = st.button("➕", key="upload_btn_row", help="Upload Document")
with icon_col2:
    refresh_clicked = st.button("🔄", key="refresh_btn_row", help="New Chat")
st.markdown('</div>', unsafe_allow_html=True)

if upload_clicked:
    st.session_state.show_uploader = True
    st.rerun()
if refresh_clicked:
    st.session_state.messages = []
    st.session_state.stage = 'idle'
    st.session_state.ppt_path = None
    st.session_state.topic = None
    st.session_state.file_content = None
    st.session_state.file_name = None
    st.session_state.parsed_slides = None
    st.session_state.awaiting_upload_confirm = False
    st.session_state.uploaded_preview = None
    st.session_state.refresh_btn_clicked = False
    st.session_state.messages.append({"role": "assistant", "content": "Chat cleared! Ready for a new topic. Type your topic or upload a document."})
    st.rerun()

# Chat input (fixed at bottom by Streamlit)
# ...existing code...

# Initialize upload confirmation state
if 'awaiting_upload_confirm' not in st.session_state:
    st.session_state.awaiting_upload_confirm = False
if 'uploaded_preview' not in st.session_state:
    st.session_state.uploaded_preview = None

# File upload section
if st.session_state.show_uploader:
    st.markdown("---")
    st.markdown("**📁 Upload Document (PDF, Word, Text, PPT, Excel, CSV)**")
    uploaded_files = document_upload_component()
    if uploaded_files:
        progress = st.progress(0, text="Processing files...")
        total = len(uploaded_files)
        for idx, (file_name, file_content) in enumerate(uploaded_files):
            if file_name in st.session_state.file_names:
                continue
            content = ""
            ext = file_name.lower().split('.')[-1]
            if ext == 'txt':
                content = file_content.decode('utf-8')
            elif ext == 'docx':
                try:
                    from docx import Document
                    import io
                    doc = Document(io.BytesIO(file_content))
                    content = '\n'.join([para.text for para in doc.paragraphs])
                except:
                    st.error(f"Could not read DOCX file: {file_name}")
            elif ext == 'pdf':
                try:
                    import PyPDF2
                    import io
                    pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                    extracted_pages = []
                    for page_num, page in enumerate(pdf_reader.pages):
                        try:
                            page_text = page.extract_text()
                            if page_text and page_text.strip():
                                extracted_pages.append(f"--- Page {page_num + 1} ---\n{page_text}")
                        except Exception as page_err:
                            print(f"[DEBUG] Could not extract page {page_num + 1}: {page_err}")
                    content = '\n\n'.join(extracted_pages)
                    print(f"[DEBUG] PDF extracted {len(extracted_pages)} pages, {len(content)} chars")
                    if not content.strip():
                        # Try alternative method with pdfplumber if PyPDF2 fails
                        try:
                            import pdfplumber
                            with pdfplumber.open(io.BytesIO(file_content)) as pdf:
                                content = '\n\n'.join([page.extract_text() or '' for page in pdf.pages])
                            print(f"[DEBUG] pdfplumber extracted {len(content)} chars")
                        except:
                            st.warning(f"⚠️ PDF text extraction limited for: {file_name}")
                except Exception as e:
                    print(f"[DEBUG] PDF extraction error: {e}")
                    st.error(f"Could not read PDF file: {file_name}")
            elif ext in ['ppt', 'pptx']:
                try:
                    from pptx import Presentation
                    import io
                    prs = Presentation(io.BytesIO(file_content))
                    slide_texts = []
                    for slide_num, slide in enumerate(prs.slides, 1):
                        slide_content = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_content.append(shape.text.strip())
                        if slide_content:
                            slide_texts.append(f"--- Slide {slide_num} ---\n" + '\n'.join(slide_content))
                    content = '\n\n'.join(slide_texts)
                    print(f"[DEBUG] PPTX extracted {len(slide_texts)} slides, {len(content)} chars")
                except Exception as e:
                    print(f"[DEBUG] PPTX extraction error: {e}")
                    st.error(f"Could not read PPTX file: {file_name}")
            elif ext in ['xlsx', 'xls', 'csv']:
                # Excel/CSV handling - store DataFrame for chart generation
                if PANDAS_AVAILABLE:
                    try:
                        import io
                        if ext == 'csv':
                            df = pd.read_csv(io.BytesIO(file_content))
                        else:
                            df = pd.read_excel(io.BytesIO(file_content))
                        # Store chart data for later use
                        st.session_state.chart_data = {'df': df, 'filename': file_name}
                        content = f"[Excel/CSV Data: {file_name} - {len(df)} rows, {len(df.columns)} columns: {', '.join(df.columns.tolist())}]"
                        print(f"[DEBUG] Excel/CSV loaded: {len(df)} rows, columns: {df.columns.tolist()}")
                    except Exception as e:
                        print(f"[DEBUG] Excel/CSV read error: {e}")
                        st.error(f"Could not read {ext.upper()} file: {file_name}")
                        content = f"[Could not read {file_name}]"
                else:
                    st.error("pandas not installed. Run: pip install pandas openpyxl")
                    content = f"[File {file_name} - pandas required]"
            else:
                content = f"[File {file_name} uploaded]"
            st.session_state.file_names.append(file_name)
            st.session_state.file_contents.append(content)
            progress.progress((idx+1)/total, text=f"Processed {idx+1}/{total} files")
        progress.empty()
        if st.session_state.file_names:
            st.session_state.show_uploader = False
            st.session_state.awaiting_upload_confirm = True
            st.rerun()

# Show success message and action buttons AFTER file upload (outside show_uploader block)
if st.session_state.awaiting_upload_confirm and st.session_state.get('file_names'):
    st.success(f"✅ {len(st.session_state.file_names)} file(s) attached successfully!")

    col1, col2 = st.columns(2)
    with col1:
        if st.button("📎 Attach more", key="attach_more_btn", use_container_width=True):
            st.session_state.show_uploader = True
            st.session_state.awaiting_upload_confirm = False
            st.rerun()
    # ─── CHART CONFIGURATION (if Excel/CSV was uploaded) ───
    chart_data = st.session_state.get('chart_data')
    if chart_data and PANDAS_AVAILABLE:
        df = chart_data['df']
        st.markdown("---")
        st.markdown(f"**📊 Chart Options for {chart_data['filename']}**")
        st.dataframe(df.head(5), use_container_width=True)

        numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
        all_cols = df.columns.tolist()

        if numeric_cols and len(all_cols) >= 2:
            chart_col1, chart_col2 = st.columns(2)
            with chart_col1:
                chart_type = st.selectbox("Chart Type", ["bar", "pie", "line"],
                                          key="chart_type_select",
                                          format_func=lambda x: {"bar": "Bar Chart", "pie": "Pie Chart", "line": "Line Chart"}[x])
            with chart_col2:
                chart_title = st.text_input("Chart Title", value=f"Data from {chart_data['filename']}",
                                            key="chart_title_input")

            st.session_state.chart_settings = {
                'chart_type': chart_type,
                'chart_title': chart_title
            }
        else:
            st.warning("Need at least 1 label column and 1 numeric column for charts.")

    with col2:
        if st.button("🚀 Generate PPT", key="generate_ppt_btn_main", use_container_width=True, type="primary"):
            # Combine all file contents for AI analysis
            combined_content = "\n\n---\n\n".join(st.session_state['file_contents'])
            file_names_str = ", ".join(st.session_state['file_names'])

            # Debug: Check if content was extracted
            print(f"[DEBUG] File contents length: {len(combined_content)}")
            print(f"[DEBUG] File names: {file_names_str}")

            # Set stage to show we're processing
            st.session_state.awaiting_upload_confirm = False

            # Check if we have actual content
            if not combined_content.strip() or len(combined_content.strip()) < 50:
                st.error("❌ No content could be extracted from uploaded files. Please check if files have readable text.")
                st.session_state.awaiting_upload_confirm = True
            else:
                with st.spinner("🤖 AI is analyzing your documents and creating slides..."):
                    try:
                        # Use AI to analyze and create proper slide content
                        generator = MultiAIGenerator()
                        language = st.session_state.get('language', 'English')

                        # Limit content size to avoid token limits (first 8000 chars)
                        content_for_ai = combined_content[:8000] if len(combined_content) > 8000 else combined_content

                        # ─────────────────────────────────────────────────────────────
                        # 🎯 AI-POWERED SMART TITLE GENERATION
                        # ─────────────────────────────────────────────────────────────
                        st.info("🎨 AI is creating an attractive title for your presentation...")
                        smart_titles = generator.generate_smart_title(
                            content=content_for_ai,
                            language=language,
                            style=st.session_state.get('theme', 'corporate')
                        )

                        if smart_titles.get('success'):
                            print(f"[DEBUG] AI-generated titles:")
                            print(f"  Main Title: {smart_titles['main_title']}")
                            print(f"  Tagline: {smart_titles['tagline']}")
                            print(f"  Subtitle: {smart_titles['subtitle']}")
                        else:
                            print(f"[DEBUG] Title generation fallback: {smart_titles.get('error', 'Unknown error')}")

                        # ─────────────────────────────────────────────────────────────
                        # 🔍 WEB SEARCH ENRICHMENT - Get related content from Wikipedia/Web
                        # ─────────────────────────────────────────────────────────────
                        web_context = ""
                        try:
                            # Extract topic/keywords from file content for search
                            first_lines = combined_content[:500].split('\n')
                            search_topic = ""
                            for line in first_lines[:5]:
                                clean = line.strip()
                                if clean and 10 < len(clean) < 100:
                                    search_topic = clean
                                    break
                            if not search_topic:
                                # Use filename as search topic
                                search_topic = re.sub(r'\.(docx|pdf|pptx|txt)$', '', file_names_str.split(',')[0], flags=re.IGNORECASE)
                                search_topic = search_topic.replace('_', ' ').replace('-', ' ')

                            print(f"[DEBUG] Web search topic: {search_topic}")

                            if google_api_key and google_cse_id and search_topic:
                                from web_search import search_google
                                results = search_google(search_topic, google_api_key, google_cse_id, num_results=3)
                                if results:
                                    web_snippets = []
                                    for r in results:
                                        snippet = r.get('snippet', '')
                                        title = r.get('title', '')
                                        if snippet:
                                            web_snippets.append(f"• {title}: {snippet}")
                                    if web_snippets:
                                        web_context = "\n\nWEB RESEARCH CONTEXT (Additional information from web search):\n" + '\n'.join(web_snippets)
                                        print(f"[DEBUG] Added web context: {len(web_context)} chars")
                        except Exception as web_err:
                            print(f"[DEBUG] Web search error (non-critical): {web_err}")

                        # Pass combined file content as user-provided content for AI to analyze
                        custom_instructions = f"""USER PROVIDED DOCUMENT CONTENT (Analyze this content and create a well-structured presentation):

{content_for_ai}
{web_context}

STRICT RULES:
- Each slide title: max 35 characters
- Each bullet: exactly 1 COMPLETE sentence, max 15 words. Never leave incomplete.
- Never truncate or use "..." — better to write less than leave incomplete.
- 4 bullet points per slide (no more, no less)
- Extract key points from the document content above
- Use web research context to enrich the information
- If Language is Hindi: Write complete short Hindi sentences, not keywords.
- Language: {language}
- Tone: Professional / Informative
- Create 8-10 slides covering all important topics"""

                        # Calculate dynamic slide count based on content size
                        content_length = len(combined_content)
                        if content_length > 50000:  # Very large document (200+ pages)
                            min_slides_dynamic = 15
                            max_slides_dynamic = 25
                        elif content_length > 20000:  # Large document (50-200 pages)
                            min_slides_dynamic = 12
                            max_slides_dynamic = 20
                        elif content_length > 10000:  # Medium document
                            min_slides_dynamic = 10
                            max_slides_dynamic = 15
                        else:  # Small document
                            min_slides_dynamic = 8
                            max_slides_dynamic = 12

                        print(f"[DEBUG] Content length: {content_length}, slides: {min_slides_dynamic}-{max_slides_dynamic}")

                        content_dict = generator.generate_ppt_content(
                            topic=f"Presentation from: {file_names_str}",
                            min_slides=min_slides_dynamic,
                            max_slides=max_slides_dynamic,
                            style=st.session_state.get('theme', 'corporate'),
                            audience="general",
                            custom_instructions=custom_instructions,
                            bullets_per_slide=st.session_state.get('bullets_per_slide', 4),
                            bullet_word_limit=30,
                            tone="professional",
                            required_phrases="",
                            forbidden_content=""
                        )

                        # Debug: Check AI response
                        print(f"[DEBUG] AI Response keys: {content_dict.keys() if isinstance(content_dict, dict) else 'Not a dict'}")

                        # Parse AI output into slides (same parser used for topic-based generation)
                        def parse_file_slides(ai_output):
                            slides = []
                            if not ai_output:
                                return slides
                            lines = [l.strip() for l in ai_output.split('\n') if l.strip()]
                            current_slide = {}
                            for line in lines:
                                # More flexible pattern to match various slide formats
                                m = re.match(r"^\*{0,2}Slide\s*(\d+)\s*[:\-–]\s*(.+?)\*{0,2}$", line, re.IGNORECASE)
                                if m:
                                    if current_slide:
                                        slides.append(current_slide)
                                    current_slide = {"slide_number": int(m.group(1)), "title": m.group(2).strip(), "bullets": []}
                                elif line.lower().startswith('main title:'):
                                    if current_slide:
                                        current_slide["main_title"] = line.split(':',1)[1].strip()
                                elif line.lower().startswith('tagline:'):
                                    if current_slide:
                                        current_slide["tagline"] = line.split(':',1)[1].strip()
                                elif line.lower().startswith('subtitle:'):
                                    if current_slide:
                                        current_slide["subtitle"] = line.split(':',1)[1].strip()
                                elif line.lower().startswith('presented by:'):
                                    if current_slide:
                                        current_slide["presented_by"] = line.split(':',1)[1].strip()
                                elif line.startswith('- ') or line.startswith('• ') or line.startswith('* ') or line.startswith('– '):
                                    bullet_text = line[2:].strip()
                                    if bullet_text and current_slide:
                                        current_slide.setdefault("bullets", []).append(bullet_text)
                                elif re.match(r'^\d+[\.\)]\s+', line):
                                    bullet_text = re.sub(r'^\d+[\.\)]\s+', '', line).strip()
                                    if bullet_text and current_slide:
                                        current_slide.setdefault("bullets", []).append(bullet_text)
                            if current_slide:
                                slides.append(current_slide)
                            # Inject presenter info
                            if slides:
                                first = slides[0]
                                if st.session_state.presenter_name:
                                    first["presented_by"] = st.session_state.presenter_name
                                if st.session_state.presenter_designation:
                                    if first.get("subtitle"):
                                        first["subtitle"] = (first.get("subtitle","") + f"\n{st.session_state.presenter_designation}").strip()
                                    else:
                                        first["subtitle"] = st.session_state.presenter_designation
                            return slides

                        ai_output = content_dict.get("output", "")
                        if not ai_output:
                            ai_output = content_dict.get("error", "")

                        print(f"[DEBUG] AI output length: {len(ai_output) if ai_output else 0}")
                        slides = parse_file_slides(ai_output)
                        print(f"[DEBUG] Parsed slides count: {len(slides)}")

                        # ─────────────────────────────────────────────────────────────
                        # 🎨 INJECT AI-GENERATED SMART TITLES INTO FIRST SLIDE
                        # ─────────────────────────────────────────────────────────────
                        if slides and smart_titles.get('success'):
                            first_slide = slides[0]
                            # Override with AI-generated titles
                            first_slide["main_title"] = smart_titles['main_title']
                            first_slide["tagline"] = smart_titles['tagline']
                            first_slide["subtitle"] = smart_titles['subtitle']
                            print(f"[DEBUG] Smart titles injected into first slide")

                        # Fallback if AI parsing returned empty slides - create content from documents
                        if not slides:
                            st.warning("⚠️ AI parsing incomplete, extracting content directly from documents...")

                            # Smart extraction from document content
                            all_lines = [l.strip() for l in combined_content.split('\n') if l.strip() and len(l.strip()) > 10]

                            # Try to find a good title from the first few lines
                            doc_title = None
                            for line in all_lines[:10]:
                                # Skip very long lines (likely paragraphs)
                                if len(line) < 100 and len(line) > 5:
                                    # Clean the line
                                    clean_line = re.sub(r'^[\d\.\-\*\#]+\s*', '', line).strip()
                                    if clean_line and len(clean_line) > 5:
                                        doc_title = clean_line[:80]  # Limit title length
                                        break

                            # Fallback to filename-based title
                            if not doc_title:
                                first_file = file_names_str.split(',')[0]
                                doc_title = re.sub(r'\.(docx|pdf|pptx|txt)$', '', first_file, flags=re.IGNORECASE).strip()
                                doc_title = doc_title.replace('_', ' ').replace('-', ' ').title()

                            # Use AI-generated smart titles if available, otherwise use fallback
                            if smart_titles.get('success'):
                                final_main_title = smart_titles['main_title']
                                final_tagline = smart_titles['tagline']
                                final_subtitle = smart_titles['subtitle']
                            else:
                                final_main_title = doc_title
                                final_tagline = "Professional Presentation"
                                final_subtitle = f"Based on {len(st.session_state['file_names'])} uploaded document(s)"

                            # Create title slide with AI-generated or extracted title
                            slides = [{
                                "slide_number": 1,
                                "title": "Title Slide",
                                "main_title": final_main_title,
                                "tagline": final_tagline,
                                "subtitle": final_subtitle,
                                "presented_by": st.session_state.get('presenter_name', ''),
                                "bullets": []
                            }]

                            # Extract meaningful paragraphs for content slides
                            paragraphs = [p.strip() for p in combined_content.split('\n\n') if p.strip() and len(p.strip()) > 30]

                            # Try to extract section titles from headings
                            section_titles = []
                            for line in all_lines:
                                # Look for potential headings (short lines, possibly numbered)
                                if 10 < len(line) < 80:
                                    clean = re.sub(r'^[\d\.\-\*\#]+\s*', '', line).strip()
                                    if clean and not clean.endswith(','):
                                        section_titles.append(clean)

                            # Create content slides (5-6 slides)
                            for i in range(min(5, max(1, len(paragraphs)))):
                                para = paragraphs[i] if i < len(paragraphs) else ""
                                if para:
                                    # Split paragraph into bullet points
                                    sentences = [s.strip() + '.' for s in re.split(r'[।\.\n]', para) if s.strip() and len(s.strip()) > 20][:3]
                                    if sentences:
                                        # Use extracted section title if available, else generic
                                        slide_title = section_titles[i] if i < len(section_titles) else f"Key Points - Section {i + 1}"
                                        # Limit title length
                                        if len(slide_title) > 60:
                                            slide_title = slide_title[:57] + "..."
                                        slides.append({
                                            "slide_number": i + 2,
                                            "title": slide_title,
                                            "bullets": sentences
                                        })

                            # Add conclusion slide
                            slides.append({
                                "slide_number": len(slides) + 1,
                                "title": "Summary & Conclusion",
                                "bullets": [
                                    f"Key insights from {len(st.session_state['file_names'])} document(s)",
                                    "Important points organized for presentation",
                                    "Review and customize as needed",
                                    "Thank you!"
                                ]
                            })

                        # Add brief AI summary to chat (full content in collapsible preview)
                        if ai_output and len(ai_output) > 50:
                            msg = f"📄 Analyzed {len(st.session_state['file_names'])} file(s) and created {len(slides)} slides. View full content in the preview below."
                            st.session_state.full_ai_output = ai_output  # Store for collapsible display
                        else:
                            msg = f"📄 Extracted content from {len(st.session_state['file_names'])} file(s) and created {len(slides)} slides."
                            st.session_state.full_ai_output = None
                        add_message("assistant", msg)

                        st.session_state.parsed_slides = slides
                        st.session_state.stage = 'generating'
                        st.session_state.topic = file_names_str
                        st.session_state.ppt_path = None
                        st.rerun()

                    except Exception as e:
                        print(f"[DEBUG] Exception: {str(e)}")
                        st.error(f"Error generating content: {str(e)}")
                        # Fallback to basic extraction with better content
                        st.info("Using direct content extraction as fallback...")

                        all_lines = [l.strip() for l in combined_content.split('\n') if l.strip() and len(l.strip()) > 10]
                        paragraphs = [p.strip() for p in combined_content.split('\n\n') if p.strip() and len(p.strip()) > 30]

                        # Extract title from content
                        doc_title = None
                        for line in all_lines[:10]:
                            if len(line) < 100 and len(line) > 5:
                                clean_line = re.sub(r'^[\d\.\-\*\#]+\s*', '', line).strip()
                                if clean_line and len(clean_line) > 5:
                                    doc_title = clean_line[:80]
                                    break
                        if not doc_title:
                            first_file = file_names_str.split(',')[0]
                            doc_title = re.sub(r'\.(docx|pdf|pptx|txt)$', '', first_file, flags=re.IGNORECASE).strip()
                            doc_title = doc_title.replace('_', ' ').replace('-', ' ').title()

                        slides = [{
                            "slide_number": 1,
                            "title": "Title Slide",
                            "main_title": doc_title,
                            "tagline": "Professional Presentation",
                            "subtitle": f"Based on {len(st.session_state['file_names'])} uploaded document(s)",
                            "presented_by": st.session_state.get('presenter_name', ''),
                            "bullets": []
                        }]

                        for i, para in enumerate(paragraphs[:5]):
                            if para:
                                sentences = [s.strip() + '.' for s in re.split(r'[।\.\n]', para) if s.strip() and len(s.strip()) > 20][:3]
                                if sentences:
                                    slides.append({
                                        "slide_number": i + 2,
                                        "title": f"Key Points - Section {i + 1}",
                                        "bullets": sentences
                                    })

                        st.session_state.parsed_slides = slides
                        st.session_state.stage = 'generating'
                        st.session_state.topic = file_names_str
                        st.session_state.ppt_path = None
                        st.rerun()

# Handle chat input (st.chat_input returns text on Enter, None otherwise)
if user_input:

    add_message("user", user_input)

    # ═══════════════════════════════════════════════════════════════════════════
    # THEME SELECTION STEP (after topic, before generating)
    # ═══════════════════════════════════════════════════════════════════════════
    if st.session_state.stage == 'awaiting_theme':
        choice = user_input.strip()
        theme_map = {'1': 'modern', '2': 'dark', '3': 'light',
                     'modern': 'modern', 'dark': 'dark', 'light': 'light'}
        chosen = theme_map.get(choice.lower())
        if chosen:
            st.session_state.theme = chosen
            st.session_state.stage = 'generating'
            st.session_state.parsed_slides = []
            st.session_state.pptxgenjs_code = None
            topic = st.session_state.get('pending_topic', choice)
            st.session_state.topic = topic
            add_message("assistant", f"Theme: **{chosen.capitalize()}**. Generating your presentation on **{topic}**...")
            st.rerun()
        else:
            add_message("assistant", "Please type **1** (Modern), **2** (Dark), or **3** (Light):")
            st.rerun()

    # ═══════════════════════════════════════════════════════════════════════════
    # 🔧 CHECK IF USER WANTS TO EDIT A SLIDE (when in preview mode)
    # ═══════════════════════════════════════════════════════════════════════════
    if st.session_state.get('in_preview_mode') and st.session_state.get('parsed_slides'):
        slide_num, instruction = detect_slide_edit_request(user_input)

        if slide_num:
            with st.spinner(f"✏️ Updating Slide {slide_num}..."):
                slides = st.session_state.parsed_slides
                language = st.session_state.get('language', 'English')
                topic = st.session_state.get('topic', 'Presentation')

                new_slide, error = regenerate_single_slide(slide_num, instruction, slides, topic, language)

                if new_slide and not error:
                    # Update the slide in the list
                    slides[slide_num - 1] = new_slide
                    st.session_state.parsed_slides = slides

                    # Clear PptxGenJS code so fallback python-pptx uses edited slides
                    st.session_state.pptxgenjs_code = None

                    # Regenerate PPT with updated slides
                    success, ppt_path = generate_ppt(slides, topic, st.session_state.theme)
                    if success:
                        st.session_state.ppt_path = ppt_path

                    add_message("assistant", f"✅ Slide {slide_num} has been updated!\n\n**New Title:** {new_slide.get('title', 'N/A')}\n**Points:** {len(new_slide.get('bullets', []))} bullet points\n\nCheck the preview above. You can make more changes or download the PPT.")
                else:
                    add_message("assistant", f"❌ Could not update Slide {slide_num}: {error}\n\nPlease try again with a clearer instruction.")

                st.rerun()
        else:
            # User is in preview mode but didn't specify a slide number
            add_message("assistant", "💡 To edit a specific slide, please mention the slide number.\n\nExamples:\n- \"Slide 3 mein bullets kam karo\"\n- \"Slide 5 ko Hindi mein likho\"\n- \"Add more points in Slide 2\"\n\nOr click **New** to create a fresh presentation on a different topic.")
            st.rerun()

    # Show the user's message immediately in chat
    with st.chat_message("user"):
        st.markdown(user_input)

    # Show loading spinner with context-aware message
    # Different message for preview mode vs initial generation
    spinner_message = "✨ Processing your changes..." if st.session_state.get('in_preview_mode') else "🤖 AI is creating your presentation..."

    with st.spinner(spinner_message):
        # If greeting, show welcome and ask for topic
        if is_greeting(user_input):
            name = st.session_state.presenter_name
            name_greeting = f" {name}" if name else ""
            if st.session_state.get('language', 'English') == 'Hindi':
                response = (f"नमस्ते{name_greeting}!\n\n"
                            "मैं आपकी मदद से पेशेवर PowerPoint प्रेजेंटेशन बना सकता हूँ।\n\n"
                            "कृपया एक विषय लिखें, या डॉक्युमेंट अपलोड करें, या टेक्स्ट पेस्ट करें।\n\n"
                            "उदाहरण: 'AI in Healthcare', 'Digital India', आदि।")
            else:
                response = (f"Hello{name_greeting}!\n\n"
                            "I can help you create a professional PowerPoint presentation.\n\n"
                            "Please enter a topic, upload a document, or paste your text.\n\n"
                            "Example: 'AI in Healthcare', 'Digital India', etc.")
            add_message("assistant", response)
            st.session_state.stage = 'idle'
            st.rerun()

        # If user input is a real topic/text/file, proceed to AI slide generation
        else:
            # Validate topic before proceeding
            is_valid, error_msg = is_valid_topic(user_input)
            if not is_valid:
                add_message("assistant", f"{error_msg}\n\nPlease try again with a proper topic.")
                st.session_state.stage = 'idle'
                st.rerun()

            # Ask theme selection before generating
            if st.session_state.stage != 'awaiting_theme':
                st.session_state.pending_topic = user_input
                add_message("assistant", "Choose a theme for your presentation:\n\n**1. Modern** - Clean white, navy & blue (recommended)\n**2. Dark** - Dark navy background, light text\n**3. Light** - White background, colorful accents\n\nType 1, 2, or 3:")
                st.session_state.stage = 'awaiting_theme'
                st.rerun()

            # ISSUE 8: Check if user provided substantial content (pasted text or uploaded file)
            user_provided_content = ""
            use_user_content = False

            # Check if file content is available
            if st.session_state.file_content:
                user_provided_content = st.session_state.file_content
                use_user_content = True
            # Check if user pasted substantial content (> 100 chars suggests they pasted content)
            elif len(user_input) > 100:
                user_provided_content = user_input
                use_user_content = True

            # Only do Google search if user didn't provide substantial content
            google_context = ""
            if not use_user_content:
                trusted_domains = ["wikipedia.org", ".gov", ".nic.in", ".org"]
                if google_api_key and google_cse_id:
                    try:
                        results = search_google(user_input, google_api_key, google_cse_id, num_results=5)
                        # Filter to trusted sources
                        trusted_results = []
                        for r in results:
                            url = r.get('link', '')
                            if any(domain in url for domain in trusted_domains):
                                trusted_results.append(r)
                        # Extract clean text (use snippet, title)
                        snippets = [r.get('snippet', '') for r in trusted_results]
                        titles = [r.get('title', '') for r in trusted_results]
                        google_context = "\n".join(titles + snippets)
                    except Exception as e:
                        st.warning(f"Google search failed: {e}")
                        google_context = ""

            # Pass context and language to AI generator
            generator = MultiAIGenerator()
            language = st.session_state.get('language', 'English')
            theme = st.session_state.get('theme', 'dark')

            # Build web context for AI
            if use_user_content:
                web_ctx = f"USER PROVIDED CONTENT (use as PRIMARY source):\n{user_provided_content}"
            else:
                web_ctx = google_context

            # ─────────────────────────────────────────────────────────────
            # 🎨 PRIMARY: Generate PptxGenJS JavaScript code via AI
            # ─────────────────────────────────────────────────────────────
            js_result = generator.generate_pptxgenjs_code(
                topic=user_input,
                theme=theme,
                web_context=web_ctx,
                language=language,
            )

            js_code = js_result.get('output', '')
            if js_code:
                st.session_state.pptxgenjs_code = js_code
                slides = []  # PptxGenJS handles all content — no placeholders needed
                print(f"[PPTXGENJS] AI generated JavaScript code ({len(js_code)} chars)")
            else:
                # ─────────────────────────────────────────────────────────
                # 🔄 FALLBACK: Generate text-based content
                # ─────────────────────────────────────────────────────────
                print(f"[FALLBACK] PptxGenJS generation failed: {js_result.get('error', 'unknown')}")
                st.session_state.pptxgenjs_code = None

                # Generate smart titles
                content_for_title = user_provided_content if use_user_content else user_input
                smart_titles = generator.generate_smart_title(
                    content=content_for_title, language=language, style=theme
                )

                custom_instructions = f"{web_ctx}\n\nLanguage: {language}\nTone: Government / Training\nSTRICT: Each bullet = 1 complete sentence, max 15 words. Each title max 35 chars."
                content_dict = generator.generate_ppt_content(
                    topic=user_input, min_slides=10, max_slides=15,
                    style=theme, audience="general",
                    custom_instructions=custom_instructions,
                    bullets_per_slide=4, bullet_word_limit=30,
                    tone="government/training",
                )

                def parse_slides(ai_output):
                    slides = []
                    if not ai_output:
                        return slides
                    lines = [l.strip() for l in ai_output.split('\n') if l.strip()]
                    current_slide = {}
                    for line in lines:
                        m = re.match(r"^\*{0,2}Slide\s*(\d+)\s*[:\-]\s*(.+?)\*{0,2}$", line, re.IGNORECASE)
                        if m:
                            if current_slide:
                                slides.append(current_slide)
                            current_slide = {"slide_number": int(m.group(1)), "title": m.group(2).strip(), "bullets": []}
                        elif line.lower().startswith('main title:'):
                            current_slide["main_title"] = line.split(':',1)[1].strip()
                        elif line.lower().startswith('tagline:'):
                            current_slide["tagline"] = line.split(':',1)[1].strip()
                        elif line.lower().startswith('subtitle:'):
                            current_slide["subtitle"] = line.split(':',1)[1].strip()
                        elif line.lower().startswith('presented by:'):
                            current_slide["presented_by"] = line.split(':',1)[1].strip()
                        elif line.startswith('- ') or line.startswith('• ') or line.startswith('* '):
                            bullet_text = line[2:].strip()
                            if bullet_text and current_slide:
                                current_slide.setdefault("bullets", []).append(bullet_text)
                        elif re.match(r'^\d+\.\s+', line):
                            bullet_text = re.sub(r'^\d+\.\s+', '', line).strip()
                            if bullet_text and current_slide:
                                current_slide.setdefault("bullets", []).append(bullet_text)
                    if current_slide:
                        slides.append(current_slide)
                    if slides:
                        first = slides[0]
                        if st.session_state.presenter_name:
                            first["presented_by"] = st.session_state.presenter_name
                        if st.session_state.presenter_designation:
                            if first.get("subtitle"):
                                first["subtitle"] = (first.get("subtitle","") + f"\n{st.session_state.presenter_designation}").strip()
                            else:
                                first["subtitle"] = st.session_state.presenter_designation
                    return slides

                ai_output = content_dict.get("output", "") or content_dict.get("error", "No AI output.")
                slides = parse_slides(ai_output)

                if slides and smart_titles.get('success'):
                    first_slide = slides[0]
                    first_slide["main_title"] = smart_titles['main_title']
                    first_slide["tagline"] = smart_titles['tagline']
                    first_slide["subtitle"] = smart_titles['subtitle']

            # Store for preview & generation
            st.session_state.full_ai_output = js_code[:500] + "..." if js_code and len(js_code) > 500 else None
            brief_msg = f"✨ Generated 10 slides for **{user_input}**. View preview below."
            add_message("assistant", brief_msg)
            st.session_state.topic = user_input
            st.session_state.parsed_slides = slides
            print(f"[DEBUG] Topic: {user_input}, Slides count: {len(slides)}")
            st.session_state.stage = 'generating'
            st.rerun()

# Generation process
if st.session_state.stage == 'generating':
    with st.chat_message("assistant"):
        progress_bar = st.progress(0)
        status_text = st.empty()

        try:
            import time
            topic = st.session_state.get('topic', '')
            theme = st.session_state.get('theme', 'modern')
            language = st.session_state.get('language', 'English')

            # Step 1: Generate JS code if not already done (e.g. after theme selection)
            status_text.text("🤖 Step 1/4: AI generating slide content...")
            progress_bar.progress(15)

            if not st.session_state.get('pptxgenjs_code'):
                generator = MultiAIGenerator()
                js_result = generator.generate_pptxgenjs_code(
                    topic=topic, theme=theme, language=language,
                    web_context=st.session_state.get('google_context', ''),
                )
                js_code = js_result.get('output', '')
                if js_code:
                    st.session_state.pptxgenjs_code = js_code
                else:
                    st.session_state.pptxgenjs_code = None

            ai_source = get_last_ai_source()

            # Step 2: Building layout
            status_text.text("🎨 Step 2/4: Building slide layout...")
            progress_bar.progress(50)
            time.sleep(0.2)

            # Step 3: Rendering PPT
            status_text.text("📊 Step 3/4: Rendering PowerPoint file...")
            progress_bar.progress(75)

            content = st.session_state.get('parsed_slides', [])

            # Step 4: Finalizing
            status_text.text("✅ Step 4/4: Finalizing...")
            progress_bar.progress(90)

            # Generate PPT
            success, ppt_path = generate_ppt(content, topic, theme)

            # Step 6: Finalizing
            status_text.text("✅ Step 6/6: Finalizing & preparing download...")
            progress_bar.progress(100)

            if success:
                status_text.text("🎉 Done! Your presentation is ready.")
                st.session_state.ppt_path = ppt_path
                st.session_state.stage = 'preview'
                st.session_state.ai_source = ai_source
                # Update PPT generated count
                try:
                    _stats_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), "visitor_count.json")
                    with open(_stats_file, 'r') as f:
                        stats = json.load(f)
                    stats["ppt_generated"] = stats.get("ppt_generated", 0) + 1
                    stats["last_ppt"] = datetime.now().isoformat()
                    with open(_stats_file, 'w') as f:
                        json.dump(stats, f)
                except:
                    pass
                time.sleep(0.5)
                progress_bar.empty()
                status_text.empty()
                st.rerun()
            else:
                progress_bar.empty()
                js_present = bool(st.session_state.get('pptxgenjs_code'))
                # Get actual node error from run_pptxgenjs debug
                js_code_debug = st.session_state.get('pptxgenjs_code', '')
                node_err = ""
                if js_code_debug:
                    import re as _re3
                    js_debug = _re3.sub(r'^\s*(const|let|var)\s+pptx\s*=\s*new\s+PptxGenJS[^\n]*\n?', '', js_code_debug, flags=_re3.MULTILINE)
                    import subprocess as _sp, tempfile as _tf
                    _fd, _tmp = _tf.mkstemp(suffix='.js')
                    with os.fdopen(_fd, 'w') as _f: _f.write(js_debug)
                    _wp = os.path.join(os.path.dirname(os.path.abspath(__file__)), "node_pptx", "pptx_wrapper.js")
                    _r = _sp.run(['node', _wp, '/tmp/debug_test.pptx', _tmp], capture_output=True, text=True, timeout=45, cwd=os.path.join(os.path.dirname(os.path.abspath(__file__)), "node_pptx"))
                    os.unlink(_tmp)
                    node_err = (_r.stderr or _r.stdout)[:200]
                status_text.error(f"Node error: {node_err[:150]}")
                print(f"[DEBUG NODE ERROR]: {node_err}")
                st.session_state.stage = 'awaiting_topic'
                add_message("assistant", f"⚠️ Node error: {node_err[:100]}. Type topic again to retry.")

        except Exception as e:
            progress_bar.empty()
            status_text.error(f"Error: {str(e)}")
            st.session_state.stage = 'awaiting_topic'
            add_message("assistant", f"Sorry, there was an error. Please try again with a different topic.")
            st.rerun()

# ═══════════════════════════════════════════════════════════════════════════════
# 📋 PREVIEW STAGE - Show slide preview with guided chat instructions
# ═══════════════════════════════════════════════════════════════════════════════
if st.session_state.stage == 'preview':
    slides = st.session_state.get('parsed_slides', [])
    ppt_path = st.session_state.get('ppt_path')

    # Show success message
    st.success(f"✅ Your presentation on **{st.session_state.topic}** is ready!")

    # ─────────────────────────────────────────────────────────────────────────────
    # 📄 COLLAPSIBLE AI OUTPUT SECTION
    # ─────────────────────────────────────────────────────────────────────────────
    full_ai_output = st.session_state.get('full_ai_output')
    if full_ai_output:
        with st.expander("📝 View Full AI Generated Content", expanded=False):
            st.text_area("AI Output", full_ai_output, height=300, disabled=True, label_visibility="collapsed")

    # ─────────────────────────────────────────────────────────────────────────────
    # 📊 SLIDE PREVIEW SECTION - Visual Thumbnails (3 per row)
    # ─────────────────────────────────────────────────────────────────────────────
    st.markdown("---")
    st.markdown("### 📊 Slide Preview")

    if slides:
        # Display slides in a 3-column grid
        cols_per_row = 3
        for row_start in range(0, len(slides), cols_per_row):
            cols = st.columns(cols_per_row)
            for col_idx, slide_idx in enumerate(range(row_start, min(row_start + cols_per_row, len(slides)))):
                slide = slides[slide_idx]
                slide_num = slide.get('slide_number', slide_idx + 1)
                title = slide.get('title', slide.get('main_title', f'Slide {slide_num}'))
                bullets = slide.get('bullets', [])

                with cols[col_idx]:
                    if slide_num == 1:
                        main_title = slide.get('main_title') or st.session_state.get('topic', 'Presentation')
                        tagline = slide.get('tagline', '')
                        subtitle = slide.get('subtitle', '')
                        st.markdown(f"""
                        <div style="background: linear-gradient(135deg, #6366f1, #8b5cf6); color: white; padding: 16px; border-radius: 12px; margin-bottom: 12px; min-height: 180px; box-shadow: 0 4px 12px rgba(99, 102, 241, 0.3);">
                            <div style="font-size: 10px; opacity: 0.8; margin-bottom: 8px;">SLIDE {slide_num}</div>
                            <div style="font-size: 14px; font-weight: bold; margin-bottom: 8px; line-height: 1.3;">{main_title[:50]}{'...' if len(main_title) > 50 else ''}</div>
                            {f'<div style="font-size: 11px; opacity: 0.9; margin-bottom: 4px;">{tagline[:40]}{"..." if len(tagline) > 40 else ""}</div>' if tagline else ''}
                            {f'<div style="font-size: 10px; opacity: 0.7;">{subtitle[:35]}{"..." if len(subtitle) > 35 else ""}</div>' if subtitle else ''}
                        </div>
                        """, unsafe_allow_html=True)
                    else:
                        bullet_preview = ''.join([f'<div style="font-size: 10px; color: #555; margin: 2px 0; line-height: 1.2;">• {b[:45]}{"..." if len(b) > 45 else ""}</div>' for b in bullets[:3]])
                        st.markdown(f"""
                        <div style="background: linear-gradient(135deg, #f8fafc, #f1f5f9); padding: 16px; border-radius: 12px; margin-bottom: 12px; min-height: 180px; border: 2px solid #e2e8f0; box-shadow: 0 2px 8px rgba(0,0,0,0.05);">
                            <div style="font-size: 10px; color: #6366f1; font-weight: bold; margin-bottom: 6px;">SLIDE {slide_num}</div>
                            <div style="font-size: 12px; font-weight: bold; color: #1e293b; margin-bottom: 10px; line-height: 1.3;">{title[:40]}{'...' if len(title) > 40 else ''}</div>
                            {bullet_preview if bullets else '<div style="font-size: 10px; color: #6366f1; font-style: italic;">✨ AI-generated content in downloaded PPT</div>'}
                        </div>
                        """, unsafe_allow_html=True)
    else:
        st.warning("No slides found in preview.")

    # ─────────────────────────────────────────────────────────────────────────────
    # 💬 GUIDED CHAT INSTRUCTION
    # ─────────────────────────────────────────────────────────────────────────────
    st.markdown("---")
    st.markdown("""
    <div style="background: linear-gradient(135deg, #fef3c7, #fde68a); padding: 16px; border-radius: 12px; border: 1px solid #f59e0b;">
        <h4 style="margin: 0 0 10px 0; color: #92400e;">💡 Want to make changes?</h4>
        <p style="margin: 0 0 12px 0; color: #78350f; font-size: 14px;">
            Type in the chat below. Examples:
        </p>
        <ul style="margin: 0; padding-left: 20px; color: #78350f; font-size: 13px;">
            <li>"Slide 3 mein bullets kam karo"</li>
            <li>"Slide 2 ko simple Hindi mein likho"</li>
            <li>"Add more points in Slide 4"</li>
        </ul>
    </div>
    """, unsafe_allow_html=True)

    # ─────────────────────────────────────────────────────────────────────────────
    # ⬇️ DOWNLOAD BUTTON (BOTTOM - always visible, big & prominent)
    # ─────────────────────────────────────────────────────────────────────────────
    st.markdown("---")
    col_dl, col_theme, col_new = st.columns([3, 2, 1])
    with col_dl:
        if ppt_path and os.path.exists(ppt_path):
            download_filename = os.path.basename(ppt_path)
            with open(ppt_path, "rb") as f:
                st.download_button(
                    "⬇️ Download PPT",
                    f.read(),
                    file_name=download_filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                    type="primary"
                )
    with col_theme:
        current_theme = st.session_state.get('theme', 'modern')
        theme_options = ["Modern", "Dark", "Light"]
        theme_map_dl = {"Modern": "modern", "Dark": "dark", "Light": "light"}
        theme_idx = {"modern": 0, "dark": 1, "light": 2}.get(current_theme, 0)
        new_theme_dl = st.selectbox("Change Theme", theme_options, index=theme_idx, label_visibility="collapsed")
        if theme_map_dl[new_theme_dl] != current_theme:
            st.session_state.theme = theme_map_dl[new_theme_dl]
            st.session_state.stage = 'regenerating'
            st.rerun()
    with col_new:
        if st.button("New", use_container_width=True):
            st.session_state.messages = []
            st.session_state.stage = 'idle'
            st.session_state.ppt_path = None
            st.session_state.topic = None
            st.session_state.parsed_slides = None
            st.session_state.file_names = []
            st.session_state.file_contents = []
            st.rerun()

    # Store that we're in preview mode for chat handling
    st.session_state.in_preview_mode = True

# Regeneration process (theme change → re-generate PptxGenJS code with new colors)
if st.session_state.stage == 'regenerating':
    with st.chat_message("assistant"):
        progress_bar = st.progress(0)
        status_text = st.empty()
        try:
            theme = st.session_state.get('theme', 'dark')
            topic = st.session_state.get('topic', 'Presentation')

            status_text.text(f"🎨 Regenerating with {theme} theme...")
            progress_bar.progress(15)

            # Re-generate PptxGenJS code with new theme colors
            generator = MultiAIGenerator()
            language = st.session_state.get('language', 'English')

            status_text.text("🧠 AI generating new presentation code...")
            progress_bar.progress(30)

            js_result = generator.generate_pptxgenjs_code(
                topic=topic,
                theme=theme,
                language=language,
            )

            js_code = js_result.get('output', '')
            if js_code:
                st.session_state.pptxgenjs_code = js_code
                print(f"[REGEN] PptxGenJS code regenerated for {theme} theme ({len(js_code)} chars)")
            else:
                print(f"[REGEN] PptxGenJS failed: {js_result.get('error')}, using fallback")
                st.session_state.pptxgenjs_code = None

            status_text.text("📊 Building PowerPoint file...")
            progress_bar.progress(70)

            content = st.session_state.get('parsed_slides', [])
            success, ppt_path = generate_ppt(content, topic, theme)

            status_text.text("✅ Finalizing...")
            progress_bar.progress(100)

            if success:
                st.session_state.ppt_path = ppt_path
                st.session_state.stage = 'done'
                progress_bar.empty()
                status_text.empty()
                add_message("assistant", f"**Regenerated!** New presentation with **{theme}** theme is ready.")
                st.rerun()
            else:
                progress_bar.empty()
                status_text.error("Failed to regenerate. Please try again.")
                st.session_state.stage = 'done'
        except Exception as e:
            progress_bar.empty()
            status_text.error(f"Error: {str(e)}")
            st.session_state.stage = 'done'
            st.rerun()

# Simple welcome message (clean chatbot style - no buttons)
if not st.session_state.messages and st.session_state.stage == 'idle':
    welcome_msg = """**Welcome to AI PPT Generator!** 🎉

I can help you create professional PowerPoint presentations in minutes.

**Here's how to get started:**

📤 **Upload a document** - Click the ➕ button above and upload PDF, Word, or PowerPoint files

📋 **Paste your content** - Copy and paste any text, article, or notes directly in the chat

✍️ **Enter a topic** - Just type a topic like "AI in Healthcare" or "Digital India" and I'll create a full presentation

**What would you like to do?** Just type below or use the buttons above! 👇"""
    st.session_state.messages.append({"role": "assistant", "content": welcome_msg})
    st.rerun()

# ============ VISITOR STATS + FOOTER ============
import json
stats_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), "visitor_count.json")
try:
    with open(stats_file, 'r') as f:
        stats = json.load(f)
except:
    stats = {"total_visits": 0, "ppt_generated": 0}

# Update visit count (only once per session)
if 'visit_counted' not in st.session_state:
    st.session_state.visit_counted = True
    stats["total_visits"] = stats.get("total_visits", 0) + 1
    stats["last_visit"] = datetime.now().isoformat()
    try:
        with open(stats_file, 'w') as f:
            json.dump(stats, f)
    except:
        pass

BASE_COUNT = 10000
total_users = BASE_COUNT + stats.get("total_visits", 0)
total_ppts = total_users  # PPTs Created = Users count

# Footer with user counter
st.markdown(f'''
<div class="footer-container">
    <span>Made with AI</span>
    <span class="footer-stat">Users: {total_users}</span>
    <span class="footer-stat">PPTs Created: {total_ppts}</span>
    <span>FREE PPT Generator</span>
</div>
''', unsafe_allow_html=True)

st.markdown('</div></div>', unsafe_allow_html=True)
