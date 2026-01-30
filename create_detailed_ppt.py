#!/usr/bin/env python3
"""
Detailed PPT Generator for Buddha and Women's Freedom
Creates comprehensive slides with tables, detailed content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

def add_table_to_slide(slide, rows, cols, left, top, width, height):
    """Add a table to the slide"""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table

def create_detailed_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Colors
    title_color = RGBColor(21, 101, 192)
    heading_color = RGBColor(255, 87, 34)
    body_color = RGBColor(33, 33, 33)
    accent_color = RGBColor(255, 193, 7)
    
    # ============ SLIDE 1: Title Slide ============
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Gradient background
    bg = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    fill = bg.fill
    fill.gradient()
    fill.gradient_angle = 45
    fill.gradient_stops[0].color.rgb = RGBColor(255, 183, 77)
    fill.gradient_stops[1].color.rgb = RGBColor(255, 112, 67)
    bg.line.fill.background()
    
    # Decorative circle
    circle = slide1.shapes.add_shape(9, Inches(7), Inches(-1), Inches(4), Inches(4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(255, 202, 40)
    circle.line.fill.background()
    
    # Title
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    tf.text = "महात्मा बुद्ध और स्त्री स्वतंत्रता"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = 'Noto Sans Devanagari'
    
    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(0.8))
    stf = subtitle_box.text_frame
    stf.text = "एक क्रांतिकारी दृष्टिकोण"
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sp.font.size = Pt(36)
    sp.font.color.rgb = RGBColor(255, 255, 255)
    sp.font.name = 'Noto Sans Devanagari'
    
    # ============ SLIDE 2: प्राचीन समाज में स्त्रियों की स्थिति ============
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Top bar
    top_bar = slide2.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.3))
    top_bar.fill.gradient()
    top_bar.fill.gradient_stops[0].color.rgb = RGBColor(0, 188, 212)
    top_bar.fill.gradient_stops[1].color.rgb = RGBColor(0, 150, 136)
    top_bar.line.fill.background()
    
    # Title
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title2.text_frame.text = "प्राचीन भारतीय समाज में स्त्रियों की स्थिति"
    title2.text_frame.paragraphs[0].font.size = Pt(32)
    title2.text_frame.paragraphs[0].font.bold = True
    title2.text_frame.paragraphs[0].font.color.rgb = heading_color
    title2.text_frame.paragraphs[0].font.name = 'Noto Sans Devanagari'
    
    # Table
    table = add_table_to_slide(slide2, 4, 2, Inches(1), Inches(1.5), Inches(8), Inches(3.5))
    
    # Table header
    table.cell(0, 0).text = "क्षेत्र"
    table.cell(0, 1).text = "स्त्रियों की स्थिति"
    
    # Table data
    table.cell(1, 0).text = "धार्मिक अधिकार"
    table.cell(1, 1).text = "वेदपाठ, यज्ञ और मोक्ष से वंचित"
    
    table.cell(2, 0).text = "शिक्षा"
    table.cell(2, 1).text = "आध्यात्मिक शिक्षा से बहिष्कृत"
    
    table.cell(3, 0).text = "सामाजिक स्वतंत्रता"
    table.cell(3, 1).text = "रीति-नीतियों की जंजीरों में जकड़ी"
    
    # Format table
    for row in table.rows:
        for cell in row.cells:
            cell.text_frame.paragraphs[0].font.size = Pt(16)
            cell.text_frame.paragraphs[0].font.name = 'Noto Sans Devanagari'
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    # Header row formatting
    for cell in table.rows[0].cells:
        cell.fill.fore_color.rgb = RGBColor(0, 150, 136)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # ============ SLIDE 3: बुद्ध का क्रांतिकारी निर्णय ============
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Side panel
    side = slide3.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.4), Inches(5.625))
    side.fill.gradient()
    side.fill.gradient_angle = 90
    side.fill.gradient_stops[0].color.rgb = RGBColor(76, 175, 80)
    side.fill.gradient_stops[1].color.rgb = RGBColor(139, 195, 74)
    side.line.fill.background()
    
    # Title
    title3 = slide3.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(9), Inches(0.6))
    title3.text_frame.text = "गौतम बुद्ध का ऐतिहासिक निर्णय (छठी शताब्दी ईसा पूर्व)"
    title3.text_frame.paragraphs[0].font.size = Pt(28)
    title3.text_frame.paragraphs[0].font.bold = True
    title3.text_frame.paragraphs[0].font.color.rgb = heading_color
    title3.text_frame.paragraphs[0].font.name = 'Noto Sans Devanagari'
    
    # Content box with detailed text
    content3 = slide3.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.8), Inches(4))
    tf3 = content3.text_frame
    tf3.word_wrap = True
    
    points = [
        ("महाप्रजापती गौतमी की माँग:", "बुद्ध की सौतेली माँ ने भिक्षुणी संघ की स्थापना की माँग की"),
        ("आनंद का तर्क:", "यदि स्त्रियाँ भी आत्मज्ञान प्राप्त कर सकती हैं, तो उन्हें क्यों रोका जाए?"),
        ("बुद्ध का निर्णय:", "महिला संन्यासिनी संगठन (भिक्षुणी संघ) की स्थापना"),
        ("महत्व:", "ब्राह्मणवादी समाज में यह एक अभूतपूर्व सामाजिक क्रांति थी"),
        ("बुद्ध का कथन:", '"नारी होना बाधा नहीं, अज्ञान होना बाधा है"')
    ]
    
    for i, (heading, detail) in enumerate(points):
        if i > 0:
            p = tf3.add_paragraph()
        else:
            p = tf3.paragraphs[0]
        
        p.text = f"• {heading}\n  {detail}"
        p.font.size = Pt(15)
        p.font.color.rgb = body_color
        p.font.name = 'Noto Sans Devanagari'
        p.space_after = Pt(14)
        p.level = 0
    
    # ============ SLIDE 4: अष्टगारधर्म (8 विशेष अनुशासन) ============
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Top gradient bar
    top4 = slide4.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.35))
    top4.fill.gradient()
    top4.fill.gradient_stops[0].color.rgb = RGBColor(103, 58, 183)
    top4.fill.gradient_stops[1].color.rgb = RGBColor(171, 71, 188)
    top4.line.fill.background()
    
    # Title
    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    title4.text_frame.text = "अष्टगारधर्म: नियम या भेदभाव?"
    title4.text_frame.paragraphs[0].font.size = Pt(32)
    title4.text_frame.paragraphs[0].font.bold = True
    title4.text_frame.paragraphs[0].font.color.rgb = heading_color
    title4.text_frame.paragraphs[0].font.name = 'Noto Sans Devanagari'
    
    # Two column layout
    # Left column - Rules
    left4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(4))
    tf_left4 = left4.text_frame
    tf_left4.word_wrap = True
    
    tf_left4.text = "8 विशेष नियम:"
    tf_left4.paragraphs[0].font.size = Pt(20)
    tf_left4.paragraphs[0].font.bold = True
    tf_left4.paragraphs[0].font.color.rgb = RGBColor(103, 58, 183)
    tf_left4.paragraphs[0].font.name = 'Noto Sans Devanagari'
    
    rules = [
        "वरिष्ठ भिक्षुणी कनिष्ठ भिक्षु को नमस्कार करेगी",
        "भिक्षुणियाँ भिक्षुओं से अनुमोदन लेंगी",
        "वर्षावास भिक्षुओं के पास करेंगी",
        "दीक्षा हेतु भिक्षु संघ की अनुमति आवश्यक"
    ]
    
    for rule in rules:
        p = tf_left4.add_paragraph()
        p.text = f"• {rule}"
        p.font.size = Pt(14)
        p.font.color.rgb = body_color
        p.font.name = 'Noto Sans Devanagari'
        p.space_after = Pt(10)
    
    # Right column - Context
    right4 = slide4.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(4))
    tf_right4 = right4.text_frame
    tf_right4.word_wrap = True
    
    tf_right4.text = "बुद्ध का दृष्टिकोण:"
    tf_right4.paragraphs[0].font.size = Pt(20)
    tf_right4.paragraphs[0].font.bold = True
    tf_right4.paragraphs[0].font.color.rgb = RGBColor(76, 175, 80)
    tf_right4.paragraphs[0].font.name = 'Noto Sans Devanagari'
    
    context = [
        "उस काल की सामाजिक संरचना में आवश्यक",
        "स्त्रियों की सुरक्षा के लिए",
        "संघ में अनुशासन बनाए रखने हेतु",
        "सामाजिक स्वीकार्यता के लिए",
        "पहली बार स्त्रियों को आध्यात्मिक स्वतंत्रता - क्रांतिकारी कदम!"
    ]
    
    for ctx in context:
        p = tf_right4.add_paragraph()
        p.text = f"✓ {ctx}"
        p.font.size = Pt(14)
        p.font.color.rgb = body_color
        p.font.name = 'Noto Sans Devanagari'
        p.space_after = Pt(10)
    
    # ============ SLIDE 5: प्रेम की नई परिभाषा ============
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background shape
    bg5 = slide5.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg5.fill.solid()
    bg5.fill.fore_color.rgb = RGBColor(255, 248, 225)
    bg5.line.fill.background()
    
    # Decorative element
    deco5 = slide5.shapes.add_shape(9, Inches(8), Inches(4), Inches(2), Inches(2))
    deco5.fill.solid()
    deco5.fill.fore_color.rgb = RGBColor(255, 193, 7)
    deco5.fill.transparency = 0.7
    deco5.line.fill.background()
    
    # Title
    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title5.text_frame.text = "प्रेम की पुनः परिभाषा: प्रकृति और आनंद की कथा"
    title5.text_frame.paragraphs[0].font.size = Pt(28)
    title5.text_frame.paragraphs[0].font.bold = True
    title5.text_frame.paragraphs[0].font.color.rgb = heading_color
    title5.text_frame.paragraphs[0].font.name = 'Noto Sans Devanagari'
    
    # Story content
    content5 = slide5.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.8), Inches(4))
    tf5 = content5.text_frame
    tf5.word_wrap = True
    
    story_parts = [
        ("कहानी:", "प्रकृति (ब्राह्मण कन्या) भिक्षु आनंद से प्रेम करने लगी। विवाह प्रस्ताव पर आनंद ने इनकार किया।"),
        ("प्रकृति का प्रश्न:", "बुद्ध ने पूछा - 'प्रेम क्या है, प्रकृति?'"),
        ("प्रकृति का उत्तर:", "'मोह, आकर्षण, पाने की इच्छा...'"),
        ("बुद्ध की परिभाषा:", "'सच्चा प्रेम पाने की नहीं, देने की इच्छा है। प्रेम वह है जो बंधन नहीं बनता, बल्कि मुक्ति देता है।'"),
        ("परिणाम:", "प्रकृति ने यह बात हृदय से ग्रहण की और बौद्ध अनुयायी बन गई"),
        ("अंबपाली का उदाहरण:", "गणिका से संन्यासिनी - हर स्त्री, चाहे उसकी सामाजिक स्थिति कोई भी हो, आत्मज्ञान प्राप्त कर सकती है")
    ]
    
    for i, (heading, detail) in enumerate(story_parts):
        if i > 0:
            p = tf5.add_paragraph()
        else:
            p = tf5.paragraphs[0]
        
        p.text = f"{heading}\n{detail}"
        p.font.size = Pt(14)
        p.font.color.rgb = body_color
        p.font.name = 'Noto Sans Devanagari'
        p.space_after = Pt(12)
    
    # ============ SLIDE 6: बुद्ध के विचारों का प्रभाव ============
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Side bar
    side6 = slide6.shapes.add_shape(1, Inches(9.6), Inches(0), Inches(0.4), Inches(5.625))
    side6.fill.gradient()
    side6.fill.gradient_angle = 90
    side6.fill.gradient_stops[0].color.rgb = RGBColor(244, 67, 54)
    side6.fill.gradient_stops[1].color.rgb = RGBColor(233, 30, 99)
    side6.line.fill.background()
    
    # Title
    title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title6.text_frame.text = "ब्राह्मणवादी समाज पर प्रभाव और आधुनिक प्रासंगिकता"
    title6.text_frame.paragraphs[0].font.size = Pt(28)
    title6.text_frame.paragraphs[0].font.bold = True
    title6.text_frame.paragraphs[0].font.color.rgb = heading_color
    title6.text_frame.paragraphs[0].font.name = 'Noto Sans Devanagari'
    
    # Table for impact
    table6 = add_table_to_slide(slide6, 5, 2, Inches(0.5), Inches(1.3), Inches(9), Inches(3.8))
    
    # Headers
    table6.cell(0, 0).text = "बुद्ध ने दिया"
    table6.cell(0, 1).text = "आधुनिक प्रासंगिकता"
    
    # Data
    table6.cell(1, 0).text = "ब्रह्मचर्य का अधिकार"
    table6.cell(1, 1).text = "करियर और आत्मनिर्भरता"
    
    table6.cell(2, 0).text = "ध्यान, साधना और निर्वाण का मार्ग"
    table6.cell(2, 1).text = "मनोवैज्ञानिक स्वतंत्रता"
    
    table6.cell(3, 0).text = "शिक्षा और संघ में भागीदारी"
    table6.cell(3, 1).text = "समान अवसर और शिक्षा का अधिकार"
    
    table6.cell(4, 0).text = "पूर्णतः स्वतंत्र आध्यात्मिक प्राणी"
    table6.cell(4, 1).text = "अपनी पहचान और पथ चुनने की स्वतंत्रता"
    
    # Format table
    for row in table6.rows:
        for cell in row.cells:
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].font.name = 'Noto Sans Devanagari'
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # Header formatting
    for cell in table6.rows[0].cells:
        cell.fill.fore_color.rgb = RGBColor(244, 67, 54)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # ============ SLIDE 7: निष्कर्ष ============
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Gradient background
    bg7 = slide7.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg7.fill.gradient()
    bg7.fill.gradient_angle = 135
    bg7.fill.gradient_stops[0].color.rgb = RGBColor(103, 58, 183)
    bg7.fill.gradient_stops[1].color.rgb = RGBColor(171, 71, 188)
    bg7.line.fill.background()
    
    # Title
    title7 = slide7.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(0.8))
    title7.text_frame.text = "निष्कर्ष"
    p7t = title7.text_frame.paragraphs[0]
    p7t.alignment = PP_ALIGN.CENTER
    p7t.font.size = Pt(44)
    p7t.font.bold = True
    p7t.font.color.rgb = RGBColor(255, 255, 255)
    p7t.font.name = 'Noto Sans Devanagari'
    
    # Conclusion points
    content7 = slide7.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tf7 = content7.text_frame
    tf7.word_wrap = True
    
    conclusions = [
        "महात्मा बुद्ध: केवल उपदेशक नहीं, समाज सुधारक भी",
        'स्त्री को "श्रद्धा" और "शक्ति" दोनों रूपों में देखा',
        "अष्टगारधर्म: समय का यथार्थ, लेकिन क्रांतिकारी कदम",
        "प्रकृति और अंबपाली: बुद्ध की शिक्षाओं का जीवंत प्रमाण",
        '"स्त्री या पुरुष - दोनों के लिए आत्मज्ञान समान रूप से सुलभ है"',
        '"प्रेम करुणा में बदले, मोह से मुक्त हो - यही बुद्ध का पथ है"'
    ]
    
    for i, conclusion in enumerate(conclusions):
        if i > 0:
            p = tf7.add_paragraph()
        else:
            p = tf7.paragraphs[0]
        
        p.text = f"• {conclusion}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = 'Noto Sans Devanagari'
        p.space_after = Pt(12)
    
    # ============ SLIDE 8: Thank You ============
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Colorful background
    bg8 = slide8.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg8.fill.gradient()
    bg8.fill.gradient_angle = 45
    bg8.fill.gradient_stops[0].color.rgb = RGBColor(0, 188, 212)
    bg8.fill.gradient_stops[1].color.rgb = RGBColor(0, 150, 136)
    bg8.line.fill.background()
    
    # Thank you message
    thanks = slide8.shapes.add_textbox(Inches(2), Inches(1.5), Inches(6), Inches(1.5))
    thanks.text_frame.text = "धन्यवाद!"
    pt = thanks.text_frame.paragraphs[0]
    pt.alignment = PP_ALIGN.CENTER
    pt.font.size = Pt(64)
    pt.font.bold = True
    pt.font.color.rgb = RGBColor(255, 255, 255)
    pt.font.name = 'Noto Sans Devanagari'
    
    # Closing message
    closing = slide8.shapes.add_textbox(Inches(2), Inches(3.2), Inches(6), Inches(1))
    closing.text_frame.text = "Like, Subscribe और Bell Icon दबाना मत भूलें!"
    pc = closing.text_frame.paragraphs[0]
    pc.alignment = PP_ALIGN.CENTER
    pc.font.size = Pt(24)
    pc.font.color.rgb = RGBColor(255, 255, 255)
    pc.font.name = 'Noto Sans Devanagari'
    
    # Save
    prs.save("output/buddha_women_freedom.pptx")
    print("✅ Detailed presentation created successfully!")
    print("📊 Total slides: 8")
    print("📍 Location: output/buddha_women_freedom.pptx")

if __name__ == "__main__":
    create_detailed_ppt()
