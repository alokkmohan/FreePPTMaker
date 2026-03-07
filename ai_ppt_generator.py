"""
AI-Powered PPT Generator
Uses OpenAI to structure content intelligently
"""

import os
import json
import shutil
import requests
import urllib.parse
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# ═══════════════════════════════════════════════════════════════════════════════
# 🖼️ POLLINATIONS.AI FREE IMAGE GENERATION (No API Key Needed)
# ═══════════════════════════════════════════════════════════════════════════════

def generate_slide_image(title, bullets=None, width=800, height=500, output_path=None):
    """
    Generate a relevant image for a slide using free APIs (no API key needed).
    Tries multiple sources: Pollinations.ai -> LoremFlickr -> Picsum.
    Returns the path to the downloaded image, or None on failure.
    """
    if not output_path:
        output_path = os.path.join(tempfile.gettempdir(), f"slide_img_{hash(title) & 0xFFFFFFFF}.jpg")

    # Extract 2-3 keywords from title for image search
    stop_words = {'the', 'a', 'an', 'in', 'on', 'of', 'and', 'for', 'to', 'is', 'are', 'was',
                  'with', 'by', 'at', 'from', 'its', 'this', 'that', 'how', 'what', 'why',
                  'ka', 'ki', 'ke', 'se', 'ko', 'hai', 'aur', 'mein', 'par', 'ya'}
    words = re.sub(r'[^\w\s]', '', title.lower()).split()
    keywords = [w for w in words if w not in stop_words and len(w) > 2][:3]
    keyword_str = ','.join(keywords) if keywords else 'business,presentation'

    # Method 1: Pollinations.ai (AI-generated)
    try:
        prompt = urllib.parse.quote(f"professional illustration {title[:60]}, modern clean style, no text")
        url = f"https://image.pollinations.ai/prompt/{prompt}?width={width}&height={height}&nologo=true"
        response = requests.get(url, timeout=20)
        if response.status_code == 200 and len(response.content) > 5000:
            content_type = response.headers.get('content-type', '')
            if 'image' in content_type or len(response.content) > 10000:
                with open(output_path, 'wb') as f:
                    f.write(response.content)
                print(f"[IMG] Pollinations OK: {output_path}")
                return output_path
    except Exception as e:
        print(f"[IMG] Pollinations skip: {e}")

    # Method 2: LoremFlickr (keyword-based real photos)
    try:
        url = f"https://loremflickr.com/{width}/{height}/{keyword_str}"
        response = requests.get(url, timeout=15, allow_redirects=True)
        if response.status_code == 200 and len(response.content) > 5000:
            with open(output_path, 'wb') as f:
                f.write(response.content)
            print(f"[IMG] LoremFlickr OK ({keyword_str}): {output_path}")
            return output_path
    except Exception as e:
        print(f"[IMG] LoremFlickr skip: {e}")

    # Method 3: Picsum (random high-quality photos as fallback)
    try:
        url = f"https://picsum.photos/{width}/{height}"
        response = requests.get(url, timeout=15, allow_redirects=True)
        if response.status_code == 200 and len(response.content) > 5000:
            with open(output_path, 'wb') as f:
                f.write(response.content)
            print(f"[IMG] Picsum OK: {output_path}")
            return output_path
    except Exception as e:
        print(f"[IMG] Picsum skip: {e}")

    print(f"[IMG] All sources failed for: {title[:40]}")
    return None


# ═══════════════════════════════════════════════════════════════════════════════
# 📊 CHART GENERATION FROM DATA (Excel/CSV)
# ═══════════════════════════════════════════════════════════════════════════════

def create_chart_image(df, chart_type="bar", title="Chart", output_path=None):
    """
    Create a chart image from a pandas DataFrame.
    Returns the path to the saved chart image, or None on failure.
    """
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt

        fig, ax = plt.subplots(figsize=(8, 5))

        # Use first column as labels, rest as data
        if len(df.columns) >= 2:
            label_col = df.columns[0]
            data_cols = df.columns[1:]

            # Limit rows for readability
            plot_df = df.head(15)

            if chart_type == "pie":
                # Pie chart uses first data column only
                values = plot_df[data_cols[0]].astype(float)
                labels = plot_df[label_col].astype(str)
                ax.pie(values, labels=labels, autopct='%1.1f%%', startangle=90,
                       colors=plt.cm.Set3.colors[:len(values)])
                ax.set_title(title, fontsize=14, fontweight='bold', pad=20)
            elif chart_type == "line":
                for col in data_cols:
                    ax.plot(plot_df[label_col].astype(str), plot_df[col].astype(float),
                            marker='o', linewidth=2, label=col)
                ax.set_xlabel(label_col, fontsize=11)
                ax.set_title(title, fontsize=14, fontweight='bold')
                ax.legend()
                plt.xticks(rotation=45, ha='right')
            else:  # bar chart (default)
                x_pos = range(len(plot_df))
                bar_width = 0.8 / len(data_cols)
                for i, col in enumerate(data_cols):
                    offset = (i - len(data_cols)/2 + 0.5) * bar_width
                    ax.bar([p + offset for p in x_pos], plot_df[col].astype(float),
                           width=bar_width, label=col)
                ax.set_xticks(x_pos)
                ax.set_xticklabels(plot_df[label_col].astype(str), rotation=45, ha='right')
                ax.set_title(title, fontsize=14, fontweight='bold')
                ax.legend()

        ax.grid(axis='y', alpha=0.3)
        plt.tight_layout()

        if not output_path:
            output_path = os.path.join(tempfile.gettempdir(), f"chart_{hash(title) & 0xFFFFFFFF}.png")
        fig.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
        plt.close(fig)
        print(f"[CHART] Chart saved: {output_path}")
        return output_path
    except Exception as e:
        print(f"[CHART] Error creating chart: {e}")
        return None

# Try to import text processor for smart truncation
try:
    from text_processor import smart_truncate, fix_broken_words, clean_bullet_point
    TEXT_PROCESSOR_AVAILABLE = True
except ImportError:
    TEXT_PROCESSOR_AVAILABLE = False
    def smart_truncate(text, max_len, suffix="..."):
        """Fallback smart truncation at word boundary"""
        if not text or len(text) <= max_len:
            return text
        truncated = text[:max_len]
        last_space = truncated.rfind(' ')
        if last_space > max_len * 0.7:
            return truncated[:last_space].rstrip('.,;:-') + suffix
        return truncated.rstrip('.,;:-') + suffix
    def fix_broken_words(text):
        return text if text else ""
    def clean_bullet_point(text):
        return text.strip() if text else ""

# ISSUE 5: Clean markdown symbols from text
import re
def clean_markdown(text):
    """Remove markdown formatting symbols from text and fix spacing"""
    if not text:
        return ""
    # Remove bold markers: **text** or __text__
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'__(.+?)__', r'\1', text)
    # Remove italic markers: *text* or _text_
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'_(.+?)_', r'\1', text)
    # Remove header markers: # ## ### etc.
    text = re.sub(r'^#{1,6}\s*', '', text, flags=re.MULTILINE)
    # Remove inline code: `text`
    text = re.sub(r'`(.+?)`', r'\1', text)
    # Remove links: [text](url) -> text
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
    # Remove bullet markers at start (will be added back as proper bullets)
    text = re.sub(r'^[\-\*•]\s*', '', text.strip())

    # ── FIX SPACING ISSUES ──
    # Add space after period/comma if missing (e.g., "word.Next" -> "word. Next")
    text = re.sub(r'([a-zA-Z\u0900-\u097F])\.([A-Z\u0900-\u097F])', r'\1. \2', text)
    text = re.sub(r'([a-zA-Z\u0900-\u097F]),([a-zA-Z\u0900-\u097F])', r'\1, \2', text)
    # Add space after colon/semicolon if missing
    text = re.sub(r'([a-zA-Z\u0900-\u097F]):([a-zA-Z\u0900-\u097F])', r'\1: \2', text)
    text = re.sub(r'([a-zA-Z\u0900-\u097F]);([a-zA-Z\u0900-\u097F])', r'\1; \2', text)
    # Fix missing space between lowercase and uppercase (camelCase splits)
    # e.g., "healthcareArtificial" -> "healthcare Artificial"
    text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)
    # Add space after closing paren if missing before a letter
    text = re.sub(r'\)([a-zA-Z\u0900-\u097F])', r') \1', text)
    # Add space before opening paren if missing after a letter
    text = re.sub(r'([a-zA-Z\u0900-\u097F])\(', r'\1 (', text)
    # Fix multiple spaces
    text = re.sub(r'\s+', ' ', text)
    # Fix space before punctuation
    text = text.replace(' ,', ',').replace(' .', '.').replace(' :', ':').replace(' ;', ';')

    return text.strip()

# Try to import image generator
try:
    from image_generator import get_slide_image, download_image, search_image
    IMAGE_GENERATION_AVAILABLE = True
except ImportError:
    IMAGE_GENERATION_AVAILABLE = False



def structure_content_with_ai(script_text, user_instructions="", min_slides=10, max_slides=20):
    """Use Ollama to structure content for slides (Ollama only)"""
    # Call Ollama API for structuring content (replace this with your actual Ollama call)
    result = ollama_structure_content(script_text, user_instructions, min_slides, max_slides)
    return result

def structure_content_basic(script_text):
    """Basic fallback structuring - flexible slide count based on content"""
    lines = [l.strip() for l in script_text.split('\n') if l.strip()]
    
    if not lines:
        return {
            "title": "Presentation",
            "subtitle": "Generated Content",
            "slides": []
        }
    
    # Extract title and subtitle
    title = lines[0] if lines else "Presentation"
    subtitle = lines[1] if len(lines) > 1 and len(lines[1]) < 100 else "Key Insights"
    
    slides = []
    current_bullets = []
    current_title = None
    
    # Process all content dynamically
    for i, line in enumerate(lines[2:] if len(lines) > 2 else lines):
        # Check if line looks like a heading/section (has colon and shortish)
        is_heading = ':' in line and len(line) < 100
        
        if is_heading:
            # Save previous slide if has enough content (minimum 4 bullets)
            if len(current_bullets) >= 4 and current_title:
                slides.append({
                    "type": "content",
                    "title": current_title,
                    "bullets": current_bullets[:6]
                })
                current_bullets = []
            elif current_bullets and current_title:
                # Store temporarily, will combine with next section
                pass
            
            # Extract new section title
            current_title = line.split(':')[0].strip()
            
            # Content after colon
            after_colon = line.split(':', 1)[1].strip() if ':' in line else ""
            if after_colon and len(after_colon) > 15:
                current_bullets.append(after_colon)
        else:
            # Regular content line
            if len(line) > 15:  # Skip very short lines
                # Set default title if none
                if not current_title:
                    current_title = "Overview"
                
                # Split very long lines into multiple bullets
                if len(line) > 200:
                    # Try to split by sentences
                    parts = line.replace('. ', '.|').split('|')
                    for part in parts:
                        if part.strip() and len(part.strip()) > 20:
                            current_bullets.append(part.strip()[:150])
                else:
                    current_bullets.append(line[:150])
                
                # Create slide when we have 5-6 bullets (good amount)
                if len(current_bullets) >= 5:
                    slides.append({
                        "type": "content",
                        "title": current_title,
                        "bullets": current_bullets[:6]
                    })
                    current_bullets = []
    
    # Add any remaining bullets as final slide(s) - ensure minimum 4 bullets
    if current_bullets:
        # Pad with generic points if less than 4
        while len(current_bullets) < 4:
            current_bullets.append(f"Important aspect of {current_title or 'this topic'}")
        
        slides.append({
            "type": "content",
            "title": current_title or "Summary",
            "bullets": current_bullets[:6]
        })
    
    return {
        "title": title,
        "subtitle": subtitle,
        "slides": slides
    }

class ModernPPTDesigner:
    """Creates beautiful, modern PPT designs"""
    
    # Professional Color Schemes
    COLOR_SCHEMES = {
        "ocean": {
            "primary": RGBColor(13, 71, 161),
            "secondary": RGBColor(3, 169, 244),
            "accent": RGBColor(255, 193, 7),
            "text": RGBColor(33, 33, 33),
            "bg": RGBColor(250, 250, 250),
            "card_bg": RGBColor(235, 240, 248),
            "card_bg_alt": RGBColor(225, 235, 245),
        },
        "forest": {
            "primary": RGBColor(27, 94, 32),
            "secondary": RGBColor(76, 175, 80),
            "accent": RGBColor(255, 152, 0),
            "text": RGBColor(33, 33, 33),
            "bg": RGBColor(250, 250, 250),
            "card_bg": RGBColor(232, 245, 233),
            "card_bg_alt": RGBColor(220, 237, 222),
        },
        "sunset": {
            "primary": RGBColor(183, 28, 28),
            "secondary": RGBColor(244, 67, 54),
            "accent": RGBColor(255, 235, 59),
            "text": RGBColor(33, 33, 33),
            "bg": RGBColor(250, 250, 250),
            "card_bg": RGBColor(255, 235, 238),
            "card_bg_alt": RGBColor(248, 225, 228),
        },
        "corporate": {
            "primary": RGBColor(26, 35, 126),
            "secondary": RGBColor(92, 107, 192),
            "accent": RGBColor(0, 188, 212),
            "text": RGBColor(33, 33, 33),
            "bg": RGBColor(255, 255, 255),
            "card_bg": RGBColor(232, 234, 246),
            "card_bg_alt": RGBColor(220, 225, 240),
        },
        "dark": {
            "primary": RGBColor(13, 27, 42),        # Dark navy
            "secondary": RGBColor(46, 196, 182),     # Teal
            "accent": RGBColor(230, 57, 70),         # Red
            "gold": RGBColor(255, 209, 102),         # Gold
            "text": RGBColor(224, 225, 221),         # Light text
            "bg": RGBColor(13, 27, 42),              # Dark navy
            "card_bg": RGBColor(27, 46, 66),         # Lighter navy
            "card_bg_alt": RGBColor(34, 58, 82),     # Alt card bg
            "muted_text": RGBColor(141, 153, 174),   # Muted
        },
        "eg": {
            "primary": RGBColor(255, 47, 47),
            "secondary": RGBColor(200, 30, 30),
            "accent": RGBColor(255, 100, 100),
            "text": RGBColor(38, 38, 38),
            "bg": RGBColor(255, 255, 255),
            "card_bg": RGBColor(255, 240, 240),
            "card_bg_alt": RGBColor(248, 230, 230),
        },
    }

    # Map UI theme names to color schemes
    THEME_ALIASES = {
        "modern": "ocean",
        "creative": "sunset",
    }

    def __init__(self, scheme="corporate"):
        resolved = self.THEME_ALIASES.get(scheme, scheme)
        self.scheme_name = resolved
        self.is_dark = (resolved == "dark")
        self.colors = self.COLOR_SCHEMES.get(resolved, self.COLOR_SCHEMES["corporate"])
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
        # Accent colors for cycling through cards
        self._accent_colors = [
            self.colors["secondary"],
            self.colors.get("accent", self.colors["secondary"]),
            self.colors.get("gold", self.colors.get("accent", self.colors["secondary"])),
        ]
    
    # ─────────────────────────────────────────────────────────────────────
    # 🔧 HELPER METHODS
    # ─────────────────────────────────────────────────────────────────────

    def _add_bg(self, slide):
        """Add full-slide background."""
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            Inches(10), Inches(5.625)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors["bg"]
        bg.line.fill.background()
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)

    def _add_card(self, slide, left, top, width, height, fill_color=None):
        """Add a card shape. Returns the shape."""
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = fill_color or self.colors.get("card_bg", RGBColor(230, 230, 240))
        card.line.fill.background()
        return card

    def _add_accent_strip(self, slide, left, top, height, color, width=0.08):
        """Add a thin colored accent strip (left edge of card)."""
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = color
        strip.line.fill.background()

    def _add_title_text(self, slide, title, y=0.3, font_size=28):
        """Add slide title text. Max 40 chars, max 32pt."""
        text_color = self.colors["text"]
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(9), Inches(0.7))
        tf = title_box.text_frame
        tf.word_wrap = True
        title_text = title[:40] if len(title) > 40 else title
        tf.text = title_text
        p = tf.paragraphs[0]
        font_size = min(font_size, 32)
        if len(title_text) > 35:
            font_size = min(font_size, 24)
        elif len(title_text) > 25:
            font_size = min(font_size, 28)
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.font.name = 'Calibri'

    def _add_red_top_bar(self, slide):
        """Add thin red accent bar at top of every slide (dark theme)."""
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.12)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.colors.get("accent", RGBColor(230, 57, 70))
        bar.line.fill.background()

    def _add_underline(self, slide, y=0.85, color=None, left=0.6, width=2.5):
        """Add thin accent underline below title."""
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(y), Inches(width), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = color or self.colors["secondary"]
        line.line.fill.background()

    def _truncate_bullet(self, text, max_words=18, max_chars=120):
        """Truncate bullet to max words/chars for visual slides."""
        text = clean_markdown(text)
        text = fix_broken_words(clean_bullet_point(text))
        words = text.split()
        if len(words) > max_words:
            text = ' '.join(words[:max_words]) + '...'
        if len(text) > max_chars:
            # Truncate at last word boundary before max_chars
            truncated = text[:max_chars]
            last_space = truncated.rfind(' ')
            if last_space > max_chars * 0.6:
                text = truncated[:last_space] + '...'
            else:
                text = truncated + '...'
        return text

    @staticmethod
    def _detect_stat_content(bullets):
        """Detect bullets with stats/numbers. Returns [(number, description)]."""
        stats = []
        for bullet in bullets:
            patterns = [
                r'(\d[\d,]*\.?\d*\s*%)',
                r'(\$[\d,]+\.?\d*\s*[BMKbmk]?)',
                r'(Rs\.?\s*[\d,]+\.?\d*\s*(?:Crore|Lakh|cr|lakh)?)',
                r'(\d[\d,]*\.?\d*\s*[xX]\b)',
                r'(\d[\d,]*\+)',
                r'(\b\d{3,}[\d,]*\b)',
            ]
            for pat in patterns:
                match = re.search(pat, bullet)
                if match:
                    num_str = match.group(1).strip()
                    desc = bullet.replace(match.group(0), '').strip(' :-,.')
                    if desc:
                        stats.append((num_str, desc))
                    break
        return stats

    # ─────────────────────────────────────────────────────────────────────
    # 🃏 CARD SLIDE - Content in visual containers (2x2 or 3-col)
    # ─────────────────────────────────────────────────────────────────────

    # Emoji icons to use in cards (cycling)
    CARD_EMOJIS = ['🔹', '🔸', '▪️', '🔷', '🔶', '⬥']

    def create_card_slide(self, title, bullets, image_path=None):
        """Content in card containers with emoji icons instead of plain bullets."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        self._add_red_top_bar(slide)
        self._add_title_text(slide, title, y=0.2)
        self._add_underline(slide, y=0.75)

        bullets = bullets[:6]  # Max 6 cards
        n = len(bullets)

        if n <= 3:
            # 3 horizontal cards
            card_w = 2.8
            card_h = 3.5
            gap = 0.3
            start_x = 0.5
            y = 1.0
            for i, bullet in enumerate(bullets):
                x = start_x + i * (card_w + gap)
                color = self._accent_colors[i % len(self._accent_colors)]
                self._add_card(slide, x, y, card_w, card_h, self.colors.get("card_bg"))
                self._add_accent_strip(slide, x, y, card_h, color)
                emoji = self.CARD_EMOJIS[i % len(self.CARD_EMOJIS)]
                tb = slide.shapes.add_textbox(
                    Inches(x + 0.2), Inches(y + 0.15),
                    Inches(card_w - 0.35), Inches(card_h - 0.3)
                )
                tf = tb.text_frame
                tf.word_wrap = True
                tf.text = f"{emoji} {self._truncate_bullet(bullet, 18, 120)}"
                p = tf.paragraphs[0]
                p.font.size = Pt(13)
                p.font.color.rgb = self.colors["text"]
                p.font.name = 'Calibri'
        else:
            # 2x2 or 2x3 grid
            card_w = 4.3
            card_h = 1.25
            gap_x = 0.4
            gap_y = 0.15
            for i, bullet in enumerate(bullets[:6]):
                col = i % 2
                row = i // 2
                x = 0.3 + col * (card_w + gap_x)
                y = 1.0 + row * (card_h + gap_y)
                color = self._accent_colors[i % len(self._accent_colors)]
                self._add_card(slide, x, y, card_w, card_h, self.colors.get("card_bg"))
                self._add_accent_strip(slide, x, y, card_h, color)
                emoji = self.CARD_EMOJIS[i % len(self.CARD_EMOJIS)]
                # Title line (first few words bold)
                text = self._truncate_bullet(bullet, 15, 120)
                words = text.split()
                card_title = ' '.join(words[:4])
                card_desc = ' '.join(words[4:]) if len(words) > 4 else ''
                # Title
                tb_title = slide.shapes.add_textbox(
                    Inches(x + 0.15), Inches(y + 0.1),
                    Inches(card_w - 0.3), Inches(0.4)
                )
                tf = tb_title.text_frame
                tf.word_wrap = True
                tf.text = f"{emoji} {card_title}"
                p = tf.paragraphs[0]
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = color
                p.font.name = 'Calibri'
                # Description
                if card_desc:
                    tb_desc = slide.shapes.add_textbox(
                        Inches(x + 0.15), Inches(y + 0.5),
                        Inches(card_w - 0.3), Inches(0.65)
                    )
                    tf2 = tb_desc.text_frame
                    tf2.word_wrap = True
                    tf2.text = card_desc
                    p2 = tf2.paragraphs[0]
                    p2.font.size = Pt(11)
                    p2.font.color.rgb = self.colors.get("muted_text", self.colors["text"])
                    p2.font.name = 'Calibri'
        return slide

    # ─────────────────────────────────────────────────────────────────────
    # 📊 STAT SLIDE - Big number callouts
    # ─────────────────────────────────────────────────────────────────────

    def create_stat_slide(self, title, stats):
        """Big stat/number callout slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        self._add_red_top_bar(slide)
        self._add_title_text(slide, title, y=0.2, font_size=26)
        self._add_underline(slide, y=0.75)

        stats = stats[:4]  # Max 4 stats
        n = len(stats)
        stat_width = 8.0 / n
        start_x = 1.0

        for i, (num_str, desc) in enumerate(stats):
            x = start_x + i * stat_width
            color = self._accent_colors[i % len(self._accent_colors)]

            # Card background
            self._add_card(slide, x - 0.2, 1.2, stat_width - 0.2, 3.5,
                          self.colors.get("card_bg"))

            # Big number
            num_box = slide.shapes.add_textbox(
                Inches(x), Inches(1.5), Inches(stat_width - 0.6), Inches(1.2)
            )
            tf = num_box.text_frame
            tf.text = num_str
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = color
            p.font.name = 'Calibri'

            # Description
            desc_box = slide.shapes.add_textbox(
                Inches(x), Inches(2.8), Inches(stat_width - 0.6), Inches(1.5)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            desc_text = desc[:60] if len(desc) > 60 else desc
            tf.text = desc_text
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(13)
            p.font.color.rgb = self.colors.get("muted_text", self.colors["text"])
            p.font.name = 'Calibri'

        return slide

    # ─────────────────────────────────────────────────────────────────────
    # ⏱️ TIMELINE SLIDE - Horizontal steps
    # ─────────────────────────────────────────────────────────────────────

    def create_timeline_slide(self, title, bullets):
        """Horizontal timeline with circles and text."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        self._add_red_top_bar(slide)
        self._add_title_text(slide, title, y=0.2, font_size=26)
        self._add_underline(slide, y=0.75)

        bullets = bullets[:5]  # Max 5 timeline items
        n = len(bullets)
        if n < 2:
            return self.create_card_slide(title, bullets)

        # Timeline line
        line_y = 2.2
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(line_y + 0.2),
            Inches(8.4), Inches(0.06)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors["secondary"]
        line.line.fill.background()

        # Circles and text
        spacing = 8.0 / (n - 1) if n > 1 else 4.0
        for i, bullet in enumerate(bullets):
            cx = 1.0 + i * spacing
            color = self._accent_colors[i % len(self._accent_colors)]

            # Circle node
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(cx - 0.25), Inches(line_y),
                Inches(0.5), Inches(0.5)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = color
            circle.line.fill.background()

            # Step number inside circle
            num_box = slide.shapes.add_textbox(
                Inches(cx - 0.25), Inches(line_y + 0.05),
                Inches(0.5), Inches(0.4)
            )
            tf = num_box.text_frame
            tf.text = str(i + 1)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = 'Calibri'

            # Text below
            text_w = min(spacing, 2.0)
            tb = slide.shapes.add_textbox(
                Inches(cx - text_w / 2), Inches(line_y + 0.65),
                Inches(text_w), Inches(2.2)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            tf.text = self._truncate_bullet(bullet, 12)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(11)
            p.font.color.rgb = self.colors["text"]
            p.font.name = 'Calibri'

        return slide

    # ─────────────────────────────────────────────────────────────────────
    # 📰 TWO-COLUMN SLIDE
    # ─────────────────────────────────────────────────────────────────────

    def create_two_column_slide(self, title, bullets, image_path=None):
        """Split content into two side-by-side containers."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        self._add_red_top_bar(slide)
        self._add_title_text(slide, title, y=0.2)
        self._add_underline(slide, y=0.75)

        mid = len(bullets) // 2
        left_bullets = bullets[:mid] if mid > 0 else bullets[:2]
        right_bullets = bullets[mid:] if mid > 0 else bullets[2:]

        # Left card
        self._add_card(slide, 0.5, 1.1, 4.3, 4.0, self.colors.get("card_bg"))
        self._add_accent_strip(slide, 0.5, 1.1, 0.06, self.colors["secondary"], width=4.3)

        tb_left = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(3.8), Inches(3.6))
        tf = tb_left.text_frame
        tf.word_wrap = True
        for i, b in enumerate(left_bullets[:3]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "  " + self._truncate_bullet(b, 15)
            p.font.size = Pt(13)
            p.font.color.rgb = self.colors["text"]
            p.font.name = 'Calibri'
            p.space_after = Pt(10)

        # Right card - use image if available, else second column
        if image_path and os.path.exists(image_path):
            try:
                slide.shapes.add_picture(
                    image_path, Inches(5.2), Inches(1.1),
                    width=Inches(4.3), height=Inches(4.0)
                )
            except:
                self._add_card(slide, 5.2, 1.1, 4.3, 4.0, self.colors.get("card_bg_alt"))
        else:
            self._add_card(slide, 5.2, 1.1, 4.3, 4.0, self.colors.get("card_bg_alt"))
            self._add_accent_strip(slide, 5.2, 1.1, 0.06, self.colors.get("accent", self.colors["secondary"]), width=4.3)

            tb_right = slide.shapes.add_textbox(Inches(5.5), Inches(1.3), Inches(3.8), Inches(3.6))
            tf = tb_right.text_frame
            tf.word_wrap = True
            for i, b in enumerate(right_bullets[:3]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = "  " + self._truncate_bullet(b, 15)
                p.font.size = Pt(13)
                p.font.color.rgb = self.colors["text"]
                p.font.name = 'Calibri'
                p.space_after = Pt(10)

        return slide

    # ─────────────────────────────────────────────────────────────────────
    # 🖼️ IMAGE + CONTENT SLIDE - Photo takes 50%, text in container
    # ─────────────────────────────────────────────────────────────────────

    def create_image_content_slide(self, title, bullets, image_path=None):
        """Full-width card with emoji bullet points."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        self._add_red_top_bar(slide)
        self._add_title_text(slide, title, y=0.2)
        self._add_underline(slide, y=0.75)

        # Full-width card container
        card_x, card_w = 0.4, 9.2
        self._add_card(slide, card_x, 0.95, card_w, 4.3, self.colors.get("card_bg"))
        self._add_accent_strip(slide, card_x, 0.95, 4.3, self.colors["secondary"])

        tb = slide.shapes.add_textbox(
            Inches(card_x + 0.3), Inches(1.1),
            Inches(card_w - 0.5), Inches(4.0)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets[:4]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            emoji = self.CARD_EMOJIS[i % len(self.CARD_EMOJIS)]
            text = self._truncate_bullet(b, 18, 120)
            p.text = f"{emoji}  {text}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors["text"]
            p.font.name = 'Calibri'
            p.space_after = Pt(14)

        return slide

    def add_gradient_background(self, slide, style="light"):
        """Add gradient background to slide"""
        # Note: python-pptx doesn't support gradients directly
        # We'll use shapes to create a layered effect
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(5.625)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors["bg"]
        bg.line.fill.background()
        
        # Send to back
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)
    
    def create_title_slide(self, main_title, tagline=None, subtitle=None, presented_by=None):
        """Beautiful title slide with dynamic sizing to prevent overflow"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        self._add_bg(slide)
        self._add_red_top_bar(slide)

        if not self.is_dark:
            # Light theme: colored header bar below red bar
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.12), Inches(10), Inches(1.1)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors["primary"]
            shape.line.fill.background()

        main_title = main_title or "Presentation"
        # Truncate title to 40 chars max
        main_title = main_title[:40] if len(main_title) > 40 else main_title

        # Count sections for spacing
        sections = 1
        if tagline: sections += 1
        if subtitle: sections += 1
        if presented_by: sections += 1

        # Dynamic title font size — max 52pt for title slide
        if len(main_title) > 35:
            title_font = 32
            title_height = 1.0
        elif len(main_title) > 25:
            title_font = 40
            title_height = 1.1
        else:
            title_font = 48
            title_height = 1.2

        y = 1.4
        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(title_height))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = main_title
        p = tf.paragraphs[0]
        p.font.size = Pt(title_font)
        p.font.bold = True
        p.font.color.rgb = self.colors["text"]
        p.font.name = 'Calibri'
        y += title_height

        # Tagline (if present)
        if tagline:
            tag_font = 20 if sections >= 4 else 24
            tagline_text = tagline[:100] if len(tagline) > 100 else tagline
            tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
            tf_tag = tagline_box.text_frame
            tf_tag.word_wrap = True
            tf_tag.text = tagline_text
            p_tag = tf_tag.paragraphs[0]
            p_tag.font.size = Pt(tag_font)
            p_tag.font.color.rgb = self.colors["secondary"]
            p_tag.font.name = 'Calibri'
            p_tag.font.italic = True
            y += 0.6

        # Subtitle (if present)
        if subtitle:
            sub_font = 16 if sections >= 4 else 20
            subtitle_text = subtitle[:120] if len(subtitle) > 120 else subtitle
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
            tf_sub = subtitle_box.text_frame
            tf_sub.word_wrap = True
            tf_sub.text = subtitle_text
            p_sub = tf_sub.paragraphs[0]
            p_sub.font.size = Pt(sub_font)
            p_sub.font.color.rgb = self.colors["secondary"]
            p_sub.font.name = 'Calibri'
            y += 0.6

        # Presented by (if present)
        if presented_by:
            pb_font = 14 if sections >= 4 else 18
            pb_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
            tf_pb = pb_box.text_frame
            tf_pb.word_wrap = True
            tf_pb.text = f"Presented by: {presented_by}"
            p_pb = tf_pb.paragraphs[0]
            p_pb.font.size = Pt(pb_font)
            p_pb.font.color.rgb = self.colors["text"]
            p_pb.font.name = 'Calibri'

        return slide
    
    def create_section_slide(self, title):
        """Section divider slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        if self.is_dark:
            # Dark theme section
            self._add_bg(slide)
            # Card in center
            self._add_card(slide, 1.5, 1.5, 7, 2.8, self.colors.get("card_bg"))
            self._add_accent_strip(slide, 1.5, 1.5, 2.8, self.colors["secondary"])
            title_color = self.colors["text"]
        else:
            # Light theme section
            bg_bottom = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625)
            )
            bg_bottom.fill.solid()
            bg_bottom.fill.fore_color.rgb = RGBColor(250, 250, 250)
            bg_bottom.line.fill.background()
            accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(3.5)
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = self.colors["secondary"]
            accent.line.fill.background()
            accent.rotation = 0
            title_color = RGBColor(255, 255, 255)

        # Title with dynamic sizing
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        tf = title_box.text_frame
        tf.text = title[:77] + "..." if len(title) > 80 else title
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        if len(title) > 60:
            p.font.size = Pt(36)
        elif len(title) > 40:
            p.font.size = Pt(42)
        else:
            p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.font.name = 'Calibri'

        return slide
    
    def create_content_slide(self, title, bullets):
        """Modern content slide with bullets - auto-adjusts header for long titles"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Dynamic header height and font based on title length
        title_text = title[:90] if len(title) > 90 else title
        if len(title_text) > 60:
            header_height = 1.3
            title_font = 22
            title_box_height = 1.0
            title_y = 0.15
        elif len(title_text) > 45:
            header_height = 1.15
            title_font = 26
            title_box_height = 0.85
            title_y = 0.15
        elif len(title_text) > 30:
            header_height = 1.0
            title_font = 30
            title_box_height = 0.7
            title_y = 0.15
        else:
            header_height = 1.0
            title_font = 36
            title_box_height = 0.7
            title_y = 0.15

        # Header bar - dynamic height
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(header_height)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.colors["primary"]
        header.line.fill.background()

        # Title on header
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(title_y),
            Inches(9), Inches(title_box_height)
        )
        tf = title_box.text_frame
        tf.text = title_text
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.font.size = Pt(title_font)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = 'Calibri'

        # Content starts after header
        content_top = header_height + 0.15

        # Side accent bar
        side_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.3), Inches(content_top),
            Inches(0.15), Inches(5.625 - content_top - 0.3)
        )
        side_bar.fill.solid()
        side_bar.fill.fore_color.rgb = self.colors["accent"]
        side_bar.line.fill.background()

        # Content area - adjusts based on header height
        content_height = 5.625 - content_top - 0.3
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(content_top),
            Inches(8.7), Inches(content_height)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        # Calculate appropriate font size based on content
        total_chars = sum(len(bullet) for bullet in bullets)
        num_bullets = len(bullets)
        
        # Detect if text contains Hindi/Devanagari characters
        has_hindi = any(ord(char) >= 0x0900 and ord(char) <= 0x097F for bullet in bullets for char in bullet)
        
        # ISSUE 4: Dynamic font sizing with INCREASED character limits to avoid truncation
        if has_hindi:
            # Hindi text takes more space, but increased limits
            if total_chars > 500 or num_bullets > 6:
                font_size = 11
                space_after = 5
                max_chars = 120  # Increased from 80
            elif total_chars > 350 or num_bullets > 5:
                font_size = 13
                space_after = 7
                max_chars = 150  # Increased from 100
            else:
                font_size = 15
                space_after = 9
                max_chars = 180  # Increased from 120
        else:
            # English text - increased limits to show more content
            if total_chars > 800 or num_bullets > 6:
                font_size = 13
                space_after = 7
                max_chars = 220  # Increased from 150
            elif total_chars > 600 or num_bullets > 5:
                font_size = 15
                space_after = 9
                max_chars = 250  # Increased from 180
            else:
                font_size = 17
                space_after = 11
                max_chars = 280  # Increased from 200
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # ISSUE 5: Clean markdown symbols first
            bullet_text = clean_markdown(bullet)
            # Clean and fix broken words
            bullet_text = fix_broken_words(clean_bullet_point(bullet_text))

            # ISSUE 4: Smart truncate with increased limits
            if len(bullet_text) > max_chars:
                bullet_text = smart_truncate(bullet_text, max_chars)

            p.text = "●  " + bullet_text
            p.font.size = Pt(font_size)
            p.font.color.rgb = self.colors["text"]
            p.font.name = 'Calibri'
            p.space_after = Pt(space_after)

        return slide
    
    def create_content_slide_with_image(self, title, bullets, image_path=None):
        """Content slide with image on right side - auto-adjusts header"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Dynamic header sizing
        clean_title = fix_broken_words(title)
        title_text = smart_truncate(clean_title, 90, "") if len(clean_title) > 90 else clean_title

        if len(title_text) > 60:
            header_height = 1.3
            title_font = 22
            title_box_height = 1.0
        elif len(title_text) > 45:
            header_height = 1.15
            title_font = 26
            title_box_height = 0.85
        elif len(title_text) > 30:
            header_height = 1.0
            title_font = 30
            title_box_height = 0.7
        else:
            header_height = 1.0
            title_font = 36
            title_box_height = 0.7

        # Header bar
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(header_height))
        header.fill.solid()
        header.fill.fore_color.rgb = self.colors["primary"]
        header.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(title_box_height))
        tf = title_box.text_frame
        tf.text = title_text
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(title_font)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = 'Calibri'

        content_top = header_height + 0.15
        content_height = 5.625 - content_top - 0.3

        # Side bar
        side_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(content_top), Inches(0.15), Inches(content_height))
        side_bar.fill.solid()
        side_bar.fill.fore_color.rgb = self.colors["accent"]
        side_bar.line.fill.background()

        # Add image if available
        content_width = 8.7
        if image_path and os.path.exists(image_path):
            try:
                slide.shapes.add_picture(image_path, Inches(5.2), Inches(content_top), width=Inches(4.3), height=Inches(content_height))
                content_width = 4.5
            except:
                pass

        # Content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(content_top), Inches(content_width), Inches(content_height))
        tf = content_box.text_frame
        tf.word_wrap = True

        total_chars = sum(len(b) for b in bullets)
        font_size = 14 if total_chars > 600 else 16 if total_chars > 400 else 18

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            clean_text = clean_markdown(bullet)
            clean_text = fix_broken_words(clean_bullet_point(clean_text))
            p.text = "● " + (smart_truncate(clean_text, 200) if len(clean_text) > 200 else clean_text)
            p.font.size = Pt(font_size)
            p.font.color.rgb = self.colors["text"]
            p.space_after = Pt(10)

        return slide
    
    def create_chart_slide(self, title, chart_image_path):
        """Create a slide with an embedded chart image"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        self._add_red_top_bar(slide)
        self._add_title_text(slide, title, y=0.2, font_size=28)
        self._add_underline(slide, y=0.75)

        # Chart image centered
        if chart_image_path and os.path.exists(chart_image_path):
            try:
                slide.shapes.add_picture(
                    chart_image_path,
                    Inches(1), Inches(1.0),
                    width=Inches(8), height=Inches(4.3)
                )
            except Exception as e:
                print(f"[CHART] Error adding chart to slide: {e}")

        return slide

    def create_end_slide(self):
        """Beautiful thank you slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        if self.is_dark:
            # Dark theme end slide
            self._add_bg(slide)
            # Card in center
            self._add_card(slide, 2.5, 1.2, 5, 3.2, self.colors.get("card_bg"))
            # Accent line at top of card
            self._add_accent_strip(slide, 2.5, 1.2, 0.06, self.colors["secondary"], width=5)
            # Thank you text
            text_box = slide.shapes.add_textbox(Inches(2.5), Inches(2.0), Inches(5), Inches(1))
            tf = text_box.text_frame
            tf.text = "Thank You!"
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(52)
            p.font.bold = True
            p.font.color.rgb = self.colors["secondary"]
            p.font.name = 'Calibri'
            # Subtitle
            sub = slide.shapes.add_textbox(Inches(2.5), Inches(3.1), Inches(5), Inches(0.6))
            tf2 = sub.text_frame
            tf2.text = "Questions & Discussion"
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            p2.font.size = Pt(18)
            p2.font.color.rgb = self.colors.get("muted_text", self.colors["text"])
            p2.font.name = 'Calibri'
        else:
            # Light theme end slide
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = self.colors["primary"]
            bg.line.fill.background()
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(3), Inches(1), Inches(4), Inches(4)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.colors["secondary"]
            circle.line.fill.background()
            text_box = slide.shapes.add_textbox(Inches(2), Inches(2.3), Inches(6), Inches(1))
            tf = text_box.text_frame
            tf.text = "Thank You!"
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(60)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = 'Calibri'

        return slide
    
    def save(self, filepath):
        """Save presentation"""
        self.prs.save(filepath)

def generate_ppt_from_template(script_text, output_path, template_path, use_ai=True, ai_instructions="", original_topic=None, min_slides=10, max_slides=20):
    """
    Generate PPT using a template file
    
    Args:
        script_text: Input script text
        output_path: Output PPT file path
        template_path: Path to template PPTX file
        use_ai: Use AI for content structuring
        ai_instructions: Additional instructions for AI (optional)
        original_topic: Original topic title to use (optional)
        min_slides: Minimum number of slides (optional)
        max_slides: Maximum number of slides (optional)
    """
    print(f"[INFO] Using template: {template_path}")
    
    # Structure content
    if use_ai:
        print("[AI] Using AI to structure content...")
        content = structure_content_with_ai(script_text, ai_instructions, min_slides, max_slides)
    else:
        print("[INFO] Using basic structuring...")
        content = structure_content_basic(script_text)
    
    # Load template
    prs = Presentation(template_path)
    
    # Use original topic if provided
    title_to_use = original_topic if original_topic else content["title"]
    
    # Try to find title slide (usually layout 0 or a slide with title placeholder)
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # Fill in title and subtitle if placeholders exist
    if title_slide.shapes.title:
        title_slide.shapes.title.text = title_to_use
    
    # Try to add subtitle if placeholder exists
    for shape in title_slide.placeholders:
        if shape.placeholder_format.idx == 1:  # Usually subtitle placeholder
            shape.text = content.get("subtitle", "")
            break
    
    # Add content slides using template layout
    # Try to use layout 1 (usually content layout) or fallback to layout 0
    content_layout_idx = 1 if len(prs.slide_layouts) > 1 else 0
    
    for slide_data in content.get("slides", []):
        slide_type = slide_data.get("type", "content")
        
        if slide_type == "section":
            # For section slides, try to use section header layout if available
            section_layout_idx = 2 if len(prs.slide_layouts) > 2 else content_layout_idx
            slide = prs.slides.add_slide(prs.slide_layouts[section_layout_idx])
            if slide.shapes.title:
                slide.shapes.title.text = slide_data["title"]
        else:
            # Regular content slide
            slide = prs.slides.add_slide(prs.slide_layouts[content_layout_idx])
            
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = slide_data["title"]
            
            # Add bullets to content placeholder
            bullets = slide_data.get("bullets", [])
            if bullets:
                # Find the content placeholder (usually idx 1)
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1 and shape.has_text_frame:
                        text_frame = shape.text_frame
                        text_frame.clear()
                        
                        for i, bullet in enumerate(bullets):
                            if i == 0:
                                p = text_frame.paragraphs[0]
                            else:
                                p = text_frame.add_paragraph()
                            p.text = bullet
                            p.level = 0
                        break
    
    # Add thank you slide using last layout
    end_layout_idx = len(prs.slide_layouts) - 1
    end_slide = prs.slides.add_slide(prs.slide_layouts[end_layout_idx])
    if end_slide.shapes.title:
        end_slide.shapes.title.text = "Thank You!"
    
    # Save
    prs.save(output_path)
    print(f"[OK] PPT created from template: {output_path}")
    
    return True

def generate_beautiful_ppt(slides_or_text, output_path, color_scheme="corporate", use_ai=True, ai_instructions="", original_topic=None, template_path=None, min_slides=10, max_slides=20, generate_ai_images=False):
    """
    Generate beautiful PPT from structured slides or script text

    Args:
        generate_ai_images: If True, generates AI images for slides using Hugging Face
    """
    print(f"DEBUG: Function called with original_topic={original_topic}, template_path={template_path}, slides={min_slides}-{max_slides}")
    print(f"DEBUG: AI Image Generation: {'ENABLED' if generate_ai_images else 'DISABLED'}")

    # If input is a list of slides (structured)
    if isinstance(slides_or_text, list) and all(isinstance(slide, dict) for slide in slides_or_text):
        slides = slides_or_text

        # Handle empty slides list
        if not slides:
            print("[WARNING] Empty slides list, creating basic presentation")
            slides = [{"title": original_topic or "Presentation", "main_title": original_topic or "Presentation"}]

        # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        # 📏 LIMIT TO EXACTLY 10 SLIDES (1 title + 8 content + 1 end)
        # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        MAX_CONTENT_SLIDES = 8
        if len(slides) > MAX_CONTENT_SLIDES + 1:  # +1 for title slide
            print(f"[LIMIT] Trimming {len(slides)} slides to {MAX_CONTENT_SLIDES + 1}")
            slides = slides[:MAX_CONTENT_SLIDES + 1]

        # Skip image generation for content slides (use emoji+card instead)
        slide_images = {}

        # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        # 📊 CREATE PPT WITH IMAGES
        # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

        designer = ModernPPTDesigner(scheme=color_scheme)

        # Title slide (use all available fields, pass separately)
        first = slides[0] if slides else {}
        main_title = first.get("main_title") or first.get("title") or original_topic or "Presentation"
        tagline = first.get("tagline") or ""
        subtitle = first.get("subtitle") or ""
        presented_by = first.get("presented_by") or ""
        designer.create_title_slide(main_title, tagline, subtitle, presented_by)

        # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        # 🎯 SMART LAYOUT ROTATION
        # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        # Layout cycle: cards -> two_column -> image_content -> cards...
        layout_cycle = ["cards", "two_column", "image_content"]
        layout_idx = 0
        # Keywords that suggest timeline/step content
        timeline_keywords = ['step', 'phase', 'stage', 'timeline', 'process',
                            'workflow', 'journey', 'roadmap', 'milestone', 'sequence']

        content_slides = slides[1:]
        for idx, slide in enumerate(content_slides, start=1):
            title = slide.get("title", "")
            bullets = slide.get("bullets", [])
            chart_image = slide.get("chart_image_path")
            title_lower = title.lower()

            if chart_image and os.path.exists(chart_image):
                # Chart slide (from Excel/CSV data)
                print(f"  [LAYOUT] Chart slide {idx}: {title[:40]}")
                designer.create_chart_slide(title, chart_image)
                continue

            if not title and not bullets:
                continue

            image_path = slide_images.get(idx, None)

            # 1. Detect stats -> stat slide
            stats = designer._detect_stat_content(bullets) if bullets else []
            if len(stats) >= 2:
                print(f"  [LAYOUT] Stat slide {idx}: {title[:40]}")
                designer.create_stat_slide(title, stats)
                continue

            # 2. Detect timeline/step content -> timeline slide
            if any(kw in title_lower for kw in timeline_keywords) and len(bullets) >= 3:
                print(f"  [LAYOUT] Timeline slide {idx}: {title[:40]}")
                designer.create_timeline_slide(title, bullets)
                continue

            # 3. Rotate through visual layouts
            current_layout = layout_cycle[layout_idx % len(layout_cycle)]
            layout_idx += 1

            if current_layout == "cards":
                print(f"  [LAYOUT] Card slide {idx}: {title[:40]}")
                designer.create_card_slide(title, bullets, image_path)
            elif current_layout == "two_column":
                print(f"  [LAYOUT] Two-column slide {idx}: {title[:40]}")
                designer.create_two_column_slide(title, bullets, image_path)
            else:  # image_content
                print(f"  [LAYOUT] Image+content slide {idx}: {title[:40]}")
                designer.create_image_content_slide(title, bullets, image_path)

        designer.create_end_slide()
        designer.save(output_path)
        print(f"[OK] Beautiful PPT created: {output_path}")
        return True
    # Otherwise, fallback to old logic
    # ...existing code...

if __name__ == "__main__":
    # Test
    sample_script = """
Uttar Pradesh's education system is undergoing major transformations.
The state is focusing on accessibility, technology, and skill development.
Operation Kayakalp aims to modernize schools.
Vocational education is being mandated in Classes 9 and 11.
    """
    
    generate_beautiful_ppt(sample_script, "test_beautiful.pptx", use_ai=False)
