#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Advanced Hindi Script to PowerPoint Converter
==============================================

Features:
- Multiple slide layouts
- Custom themes
- Image support
- Bullet points formatting
- Section dividers
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

class PPTTheme:
    """PPT color themes"""
    
    # Professional Blue Theme
    PROFESSIONAL = {
        'title': RGBColor(21, 101, 192),
        'heading': RGBColor(21, 101, 192),
        'body': RGBColor(33, 33, 33),
        'accent': RGBColor(255, 152, 0),
        'bg': RGBColor(255, 255, 255)
    }
    
    # Warm Orange Theme
    WARM = {
        'title': RGBColor(230, 81, 0),
        'heading': RGBColor(244, 67, 54),
        'body': RGBColor(33, 33, 33),
        'accent': RGBColor(255, 193, 7),
        'bg': RGBColor(255, 255, 255)
    }
    
    # Cool Green Theme
    COOL = {
        'title': RGBColor(46, 125, 50),
        'heading': RGBColor(76, 175, 80),
        'body': RGBColor(33, 33, 33),
        'accent': RGBColor(0, 150, 136),
        'bg': RGBColor(255, 255, 255)
    }
    
    # Purple Theme
    PURPLE = {
        'title': RGBColor(106, 27, 154),
        'heading': RGBColor(142, 36, 170),
        'body': RGBColor(33, 33, 33),
        'accent': RGBColor(171, 71, 188),
        'bg': RGBColor(255, 255, 255)
    }

class AdvancedPPTMaker:
    """Advanced PPT maker with multiple features"""
    
    def __init__(self, theme='PROFESSIONAL'):
        """
        Initialize PPT maker
        
        Args:
            theme: Color theme name (PROFESSIONAL, WARM, COOL, PURPLE)
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
        
        # Set theme
        theme_dict = getattr(PPTTheme, theme, PPTTheme.PROFESSIONAL)
        self.colors = theme_dict
        
        self.slide_count = 0
    
    def add_title_slide(self, title, subtitle=""):
        """
        Add title slide with gradient background
        
        Args:
            title: Main title
            subtitle: Subtitle text
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Add colorful background shape
        bg_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(10), Inches(5.625)
        )
        fill = bg_shape.fill
        fill.gradient()
        fill.gradient_angle = 45
        fill.gradient_stops[0].color.rgb = RGBColor(255, 183, 77)  # Light orange
        fill.gradient_stops[1].color.rgb = RGBColor(255, 112, 67)  # Darker orange
        bg_shape.line.fill.background()
        
        # Add decorative circle
        circle = slide.shapes.add_shape(
            9,  # Oval
            Inches(7), Inches(-1),
            Inches(4), Inches(4)
        )
        circle_fill = circle.fill
        circle_fill.solid()
        circle_fill.fore_color.rgb = RGBColor(255, 202, 40)  # Yellow
        circle.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.8), 
            Inches(9), Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
        title_frame.paragraphs[0].font.name = 'Noto Sans Devanagari'
        
        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(3.2), 
                Inches(8), Inches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            subtitle_frame.paragraphs[0].font.size = Pt(32)
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
            subtitle_frame.paragraphs[0].font.name = 'Noto Sans Devanagari'
        
        self.slide_count += 1
        return slide
    
    def add_section_slide(self, section_title):
        """
        Add section divider slide
        
        Args:
            section_title: Section title text
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Add colorful gradient background
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(10), Inches(5.625)
        )
        fill = shape.fill
        fill.gradient()
        fill.gradient_angle = 135
        fill.gradient_stops[0].color.rgb = RGBColor(103, 58, 183)  # Deep purple
        fill.gradient_stops[1].color.rgb = RGBColor(171, 71, 188)  # Light purple
        shape.line.fill.background()
        
        # Section title (white text on colored bg)
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), 
            Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].font.name = 'Noto Sans Devanagari'
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        self.slide_count += 1
        return slide
    
    def add_bullet_slide(self, title, bullets):
        """
        Add slide with bullet points
        
        Args:
            title: Slide title
            bullets: List of bullet points
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Add colorful top bar
        top_bar = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(10), Inches(0.3)
        )
        bar_fill = top_bar.fill
        bar_fill.gradient()
        bar_fill.gradient_stops[0].color.rgb = RGBColor(0, 188, 212)  # Cyan
        bar_fill.gradient_stops[1].color.rgb = RGBColor(0, 150, 136)  # Teal
        top_bar.line.fill.background()
        
        # Add colorful accent circle
        accent_circle = slide.shapes.add_shape(
            9,  # Oval
            Inches(0.2), Inches(0.35),
            Inches(0.5), Inches(0.5)
        )
        circle_fill = accent_circle.fill
        circle_fill.solid()
        circle_fill.fore_color.rgb = RGBColor(255, 193, 7)  # Yellow
        accent_circle.line.fill.background()
        
        # Title with accent bar
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), 
            Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['heading']
        title_frame.paragraphs[0].font.name = 'Noto Sans Devanagari'
        
        # Bullets
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.3), 
            Inches(8.5), Inches(3.8)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = "• " + bullet
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['body']
            p.font.name = 'Noto Sans Devanagari'
            p.space_after = Pt(12)
        
        self.slide_count += 1
        return slide
    
    def add_content_slide(self, title, content_lines):
        """
        Add slide with paragraph content
        
        Args:
            title: Slide title
            content_lines: List of content lines
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Add colorful side panel
        side_panel = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(0.4), Inches(5.625)
        )
        panel_fill = side_panel.fill
        panel_fill.gradient()
        panel_fill.gradient_angle = 90
        panel_fill.gradient_stops[0].color.rgb = RGBColor(76, 175, 80)  # Green
        panel_fill.gradient_stops[1].color.rgb = RGBColor(139, 195, 74)  # Light green
        side_panel.line.fill.background()
        
        # Add decorative shapes
        deco_shape = slide.shapes.add_shape(
            9,  # Oval
            Inches(8.5), Inches(4.5),
            Inches(1.5), Inches(1.5)
        )
        deco_fill = deco_shape.fill
        deco_fill.solid()
        deco_fill.fore_color.rgb = RGBColor(255, 235, 59)  # Yellow
        deco_fill.transparency = 0.7
        deco_shape.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), 
            Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['heading']
        title_frame.paragraphs[0].font.name = 'Noto Sans Devanagari'
        
        # Content
        content_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.3), 
            Inches(8.6), Inches(3.8)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, line in enumerate(content_lines):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = line
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['body']
            p.font.name = 'Noto Sans Devanagari'
            p.space_after = Pt(10)
        
        self.slide_count += 1
        return slide
    
    def add_two_column_slide(self, title, left_content, right_content):
        """
        Add slide with two columns
        
        Args:
            title: Slide title
            left_content: List of lines for left column
            right_content: List of lines for right column
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), 
            Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['heading']
        
        # Left column
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3), 
            Inches(4.5), Inches(3.8)
        )
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        
        for i, line in enumerate(left_content):
            if i == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['body']
        
        # Right column
        right_box = slide.shapes.add_textbox(
            Inches(5.2), Inches(1.3), 
            Inches(4.5), Inches(3.8)
        )
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        for i, line in enumerate(right_content):
            if i == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['body']
        
        self.slide_count += 1
        return slide
    
    def add_end_slide(self, message="धन्यवाद!"):
        """
        Add final thank you slide
        
        Args:
            message: End message
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # End message
        end_box = slide.shapes.add_textbox(
            Inches(2), Inches(2), 
            Inches(6), Inches(1.5)
        )
        end_frame = end_box.text_frame
        end_frame.text = message
        end_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        end_frame.paragraphs[0].font.size = Pt(54)
        end_frame.paragraphs[0].font.bold = True
        end_frame.paragraphs[0].font.color.rgb = self.colors['accent']
        end_frame.paragraphs[0].font.name = 'Noto Sans Devanagari'
        end_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        self.slide_count += 1
        return slide
    
    def save(self, filename):
        """
        Save presentation
        
        Args:
            filename: Output filename
        """
        self.prs.save(filename)
        print(f"✅ Presentation saved: {filename}")
        print(f"📊 Total slides: {self.slide_count}")

# Example usage
if __name__ == "__main__":
    # Create PPT with WARM theme
    maker = AdvancedPPTMaker(theme='PROFESSIONAL')
    
    # Title slide
    maker.add_title_slide(
        "स्वामी विवेकानंद और गौतम बुद्ध",
        "करुणा, आत्मज्ञान और सत्य की साझी चेतना"
    )
    
    # Section slide
    maker.add_section_slide("भाग १: परिचय")
    
    # Bullet slide
    maker.add_bullet_slide(
        "मुख्य बिंदु",
        [
            "करुणा और मानव सेवा",
            "आत्म-ज्ञान और ध्यान",
            "सत्य की खोज",
            "शांति और अहिंसा"
        ]
    )
    
    # Content slide
    maker.add_content_slide(
        "विवेकानंद की शिक्षाएं",
        [
            "स्वामी विवेकानंद ने बुद्ध की करुणा को अपनाया।",
            "उन्होंने मानवता की सेवा को सर्वोच्च धर्म बताया।",
            "ध्यान और आत्मनिरीक्षण पर जोर दिया।"
        ]
    )
    
    # End slide
    maker.add_end_slide()
    
    # Save
    maker.save("advanced_presentation.pptx")
