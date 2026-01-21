#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Claude-Powered Professional PPT Generator
Generates high-quality presentations using Claude API for content analysis
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
from typing import Dict, List, Optional
from claude_content_analyzer import ClaudeContentAnalyzer

# Professional color schemes
COLOR_SCHEMES = {
    "professional": {
        "primary": RGBColor(30, 58, 138),      # Deep Blue
        "secondary": RGBColor(5, 150, 105),    # Emerald Green
        "accent": RGBColor(220, 38, 38),       # Red
        "dark": RGBColor(31, 41, 55),          # Dark Gray
        "light": RGBColor(107, 114, 128),      # Light Gray
        "bg": RGBColor(255, 255, 255)          # White
    },
    "government": {
        "primary": RGBColor(30, 58, 138),
        "secondary": RGBColor(5, 150, 105),
        "accent": RGBColor(220, 38, 38),
        "dark": RGBColor(31, 41, 55),
        "light": RGBColor(107, 114, 128),
        "bg": RGBColor(255, 255, 255)
    },
    "corporate": {
        "primary": RGBColor(59, 130, 246),     # Blue
        "secondary": RGBColor(16, 185, 129),   # Green
        "accent": RGBColor(245, 158, 11),      # Orange
        "dark": RGBColor(17, 24, 39),
        "light": RGBColor(156, 163, 175),
        "bg": RGBColor(255, 255, 255)
    },
    "technical": {
        "primary": RGBColor(99, 102, 241),     # Indigo
        "secondary": RGBColor(139, 92, 246),   # Purple
        "accent": RGBColor(236, 72, 153),      # Pink
        "dark": RGBColor(17, 24, 39),
        "light": RGBColor(156, 163, 175),
        "bg": RGBColor(255, 255, 255)
    }
}


class ClaudePPTGenerator:
    """Generate professional PPTs using Claude-analyzed content"""
    
    def __init__(self, color_scheme: str = "professional"):
        """Initialize PPT generator with color scheme"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
        
        self.colors = COLOR_SCHEMES.get(color_scheme, COLOR_SCHEMES["professional"])
        self.color_scheme = color_scheme
    
    def add_title_slide(self, title: str, subtitle: str = "", presenter: str = ""):
        """Add professional title slide"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Clear default shapes
        for shape in list(slide.shapes):
            sp = shape.element
            sp.getparent().remove(sp)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9), Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors["primary"]
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.8), Inches(9), Inches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = subtitle
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = self.colors["light"]
            subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Presenter
        if presenter:
            presenter_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.5), Inches(9), Inches(0.6)
            )
            presenter_frame = presenter_box.text_frame
            presenter_para = presenter_frame.paragraphs[0]
            presenter_para.text = f"Presented By: {presenter}"
            presenter_para.font.size = Pt(16)
            presenter_para.font.color.rgb = self.colors["dark"]
            presenter_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(
        self, 
        title: str, 
        content: List[str], 
        speaker_notes: str = ""
    ):
        """Add professional content slide with bullets"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Clear default shapes
        for shape in list(slide.shapes):
            sp = shape.element
            sp.getparent().remove(sp)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors["primary"]
        
        # Content area
        content_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.2), Inches(8.8), Inches(3.8)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        # Add bullets
        for i, bullet_text in enumerate(content):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = bullet_text
            p.font.size = Pt(14)
            p.font.color.rgb = self.colors["dark"]
            p.level = 0
            p.space_before = Pt(6)
            p.space_after = Pt(6)
        
        # Add speaker notes
        if speaker_notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = speaker_notes
        
        return slide
    
    def add_section_divider(self, section_title: str):
        """Add section divider slide"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Clear default shapes
        for shape in list(slide.shapes):
            sp = shape.element
            sp.getparent().remove(sp)
        
        # Background shape
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(5.625)
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self.colors["primary"]
        bg_shape.line.fill.background()
        
        # Section title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(9), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = section_title
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def generate_from_structure(
        self, 
        ppt_structure: Dict, 
        presenter: str = "",
        add_dividers: bool = False
    ) -> bool:
        """
        Generate complete PPT from Claude-analyzed structure
        
        Args:
            ppt_structure: Dict with title, subtitle, slides
            presenter: Presenter name
            add_dividers: Add section divider slides
        
        Returns:
            True if successful
        """
        try:
            # Add title slide
            self.add_title_slide(
                title=ppt_structure.get("title", "Presentation"),
                subtitle=ppt_structure.get("subtitle", ""),
                presenter=presenter
            )
            
            # Add content slides
            slides = ppt_structure.get("slides", [])
            
            for slide_data in slides:
                title = slide_data.get("title", "")
                content = slide_data.get("content", [])
                notes = slide_data.get("speaker_notes", "")
                
                # Add divider if this is a major section
                if add_dividers and "section" in title.lower():
                    self.add_section_divider(title)
                
                # Add content slide
                self.add_content_slide(
                    title=title,
                    content=content,
                    speaker_notes=notes
                )
            
            return True
        
        except Exception as e:
            print(f"Error generating PPT: {str(e)}")
            return False
    
    def save(self, output_path: str):
        """Save presentation to file"""
        try:
            self.prs.save(output_path)
            return True
        except Exception as e:
            print(f"Error saving PPT: {str(e)}")
            return False


# High-level convenience functions

def create_ppt_from_file(
    file_path,
    output_path: str,
    style: str = "professional",
    min_slides: int = 10,
    max_slides: int = 20,
    audience: str = "general",
    presenter: str = "",
    custom_instructions: str = "",
    api_key: str = None
) -> bool:
    """
    Complete workflow: File → Claude Analysis → Professional PPT
    
    Args:
        file_path: Input file (Word/PDF/Text)
        output_path: Output PPT path
        style: Presentation style
        min_slides: Minimum slides
        max_slides: Maximum slides
        audience: Target audience
        presenter: Presenter name
        custom_instructions: Additional instructions
        api_key: Claude API key
    
    Returns:
        True if successful
    """
    try:
        print("📄 Analyzing file(s) with Claude...")
        from claude_content_analyzer import ClaudeContentAnalyzer
        analyzer = ClaudeContentAnalyzer(api_key=api_key)
        # Support both single file and list of files
        if isinstance(file_path, list):
            all_text = []
            for fp in file_path:
                print(f"[DEBUG] Processing file: {fp}")
                all_text.append(analyzer.extract_text_from_file(fp))
            combined_content = '\n\n'.join(all_text)
        else:
            print(f"[DEBUG] Processing file: {file_path}")
            combined_content = analyzer.extract_text_from_file(file_path)
        print(f"[DEBUG] Combined content length: {len(combined_content)}")
        ppt_structure = analyzer.analyze_for_ppt(
            content=combined_content,
            style=style,
            min_slides=min_slides,
            max_slides=max_slides,
            audience=audience,
            custom_instructions=custom_instructions
        )
        print(f"✅ Analysis complete: {len(ppt_structure['slides'])} slides planned")
        print(f"[DEBUG] PPT Structure: {ppt_structure}")
        print("🎨 Creating professional presentation...")
        generator = ClaudePPTGenerator(color_scheme=style)
        success = generator.generate_from_structure(
            ppt_structure=ppt_structure,
            presenter=presenter
        )
        if success:
            print("💾 Saving presentation...")
            generator.save(output_path)
            print(f"✅ Presentation saved: {output_path}")
            return True
        else:
            print("❌ Failed to generate presentation")
            return False
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        import traceback
        traceback.print_exc()
        return False


def create_ppt_from_topic(
    topic: str,
    output_path: str,
    style: str = "professional",
    min_slides: int = 10,
    max_slides: int = 20,
    audience: str = "general",
    presenter: str = "",
    custom_instructions: str = "",
    api_key: str = None
) -> bool:
    """
    Complete workflow: Topic → Claude Research → Professional PPT
    
    Args:
        topic: Presentation topic
        output_path: Output PPT path
        style: Presentation style
        min_slides: Minimum slides
        max_slides: Maximum slides
        audience: Target audience
        presenter: Presenter name
        custom_instructions: Additional instructions
        api_key: Claude API key
    
    Returns:
        True if successful
    """
    try:
        print(f"🔍 Researching topic with Claude: {topic}")
        
        # Generate content and structure with Claude
        ppt_structure = generate_ppt_from_topic(
            topic=topic,
            style=style,
            min_slides=min_slides,
            max_slides=max_slides,
            audience=audience,
            custom_instructions=custom_instructions,
            api_key=api_key
        )
        
        print(f"✅ Research complete: {len(ppt_structure['slides'])} slides planned")
        
        # Generate PPT
        print("🎨 Creating professional presentation...")
        generator = ClaudePPTGenerator(color_scheme=style)
        
        success = generator.generate_from_structure(
            ppt_structure=ppt_structure,
            presenter=presenter
        )
        
        if success:
            print("💾 Saving presentation...")
            generator.save(output_path)
            print(f"✅ Presentation saved: {output_path}")
            return True
        else:
            print("❌ Failed to generate presentation")
            return False
    
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return False


# Example usage
if __name__ == "__main__":
    # Test with a topic
    success = create_ppt_from_topic(
        topic="Artificial Intelligence in Education",
        output_path="ai_education.pptx",
        style="professional",
        min_slides=12,
        max_slides=15,
        audience="educators",
        presenter="Alok Mohan"
    )
    
    if success:
        print("🎉 Success! Professional presentation created.")
    else:
        print("❌ Failed to create presentation.")
