#!/usr/bin/env python3
"""
PowerPoint to Images Converter
Converts each slide of a PPT to individual image files
- Windows: Uses PowerPoint COM automation (best quality)
- Linux/Cloud: Uses python-pptx + PIL (simple preview)
"""

import os
import sys
from pathlib import Path

def ppt_to_images_fallback(ppt_file, output_dir="output/slides"):
    """
    Create simple preview images using python-pptx and PIL
    Works on Linux/Cloud without PowerPoint
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from PIL import Image, ImageDraw, ImageFont

        print("[PPT] Using python-pptx fallback for preview...")

        # Load presentation
        prs = Presentation(ppt_file)

        # Image settings
        width, height = 1280, 720
        bg_color = (255, 255, 255)  # White background
        title_color = (51, 51, 51)  # Dark gray for title
        text_color = (80, 80, 80)  # Gray for content

        # Try to load a font
        try:
            title_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 36)
            text_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
        except:
            try:
                title_font = ImageFont.truetype("arial.ttf", 36)
                text_font = ImageFont.truetype("arial.ttf", 24)
            except:
                title_font = ImageFont.load_default()
                text_font = ImageFont.load_default()

        slide_count = len(prs.slides)
        print(f"[PPT] Found {slide_count} slides")

        for idx, slide in enumerate(prs.slides, 1):
            # Create image
            img = Image.new('RGB', (width, height), bg_color)
            draw = ImageDraw.Draw(img)

            # Add gradient header
            for y in range(100):
                r = int(102 + (118 - 102) * y / 100)
                g = int(126 + (75 - 126) * y / 100)
                b = int(234 + (162 - 234) * y / 100)
                draw.line([(0, y), (width, y)], fill=(r, g, b))

            # Extract title
            title_text = ""
            content_texts = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        if shape.shape_type == 14 or (hasattr(shape, 'is_placeholder') and shape.placeholder_format and shape.placeholder_format.type == 1):
                            title_text = text
                        else:
                            # Split by newlines and add as bullet points
                            for line in text.split('\n'):
                                line = line.strip()
                                if line and line != title_text:
                                    content_texts.append(line)

            # If no title found, use first text
            if not title_text and content_texts:
                title_text = content_texts.pop(0)

            # Draw title (white on gradient)
            if title_text:
                # Truncate if too long
                if len(title_text) > 60:
                    title_text = title_text[:57] + "..."
                draw.text((40, 30), title_text, fill=(255, 255, 255), font=title_font)

            # Draw slide number
            draw.text((width - 80, 30), f"Slide {idx}", fill=(255, 255, 255, 180), font=text_font)

            # Draw content
            y_pos = 130
            max_lines = 12
            line_count = 0

            for text in content_texts[:max_lines]:
                if y_pos > height - 50:
                    break
                # Truncate long lines
                if len(text) > 80:
                    text = text[:77] + "..."
                # Add bullet point
                bullet_text = f"  {text}"
                draw.text((40, y_pos), bullet_text, fill=text_color, font=text_font)
                y_pos += 40
                line_count += 1

            # Add "..." if more content
            if len(content_texts) > max_lines:
                draw.text((40, y_pos), "  ...", fill=text_color, font=text_font)

            # Save image
            image_path = os.path.join(output_dir, f"slide_{idx:02d}.png")
            img.save(image_path, "PNG")
            print(f"  [OK] Slide {idx} saved")

        print("[OK] All slides converted (fallback mode)!")
        return True

    except Exception as e:
        print(f"[ERROR] Fallback conversion failed: {e}")
        return False


def ppt_to_images(ppt_file, output_dir="output/slides"):
    """
    Convert PowerPoint slides to images
    - First tries Windows COM automation (best quality)
    - Falls back to python-pptx + PIL (works everywhere)
    """

    # Create output directory
    os.makedirs(output_dir, exist_ok=True)

    print(f"[PPT] Converting PPT to images...")
    print(f"[PPT] Input: {ppt_file}")
    print(f"[PPT] Output: {output_dir}")

    # Check if file exists
    if not os.path.exists(ppt_file):
        print(f"[ERROR] File not found: {ppt_file}")
        return False

    # Get absolute paths
    ppt_file = os.path.abspath(ppt_file)
    output_dir = os.path.abspath(output_dir)

    # Check if we're on Windows
    if sys.platform == 'win32':
        try:
            # Try Windows COM automation (requires PowerPoint installed)
            import comtypes.client

            print("[PPT] Using PowerPoint COM automation...")

            # Start PowerPoint
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = 1

            # Open presentation
            presentation = powerpoint.Presentations.Open(ppt_file, WithWindow=False)

            # Export each slide as PNG
            total_slides = presentation.Slides.Count
            print(f"[PPT] Found {total_slides} slides")

            for i in range(1, total_slides + 1):
                slide = presentation.Slides(i)
                image_path = os.path.join(output_dir, f"slide_{i:02d}.png")
                slide.Export(image_path, "PNG", 1280, 720)
                print(f"  [OK] Slide {i} saved")

            # Close presentation
            presentation.Close()
            powerpoint.Quit()

            print("[OK] All slides converted!")
            return True

        except ImportError:
            print("[INFO] comtypes not available, using fallback...")
            return ppt_to_images_fallback(ppt_file, output_dir)
        except Exception as e:
            error_msg = str(e)
            if "PowerPoint.Application" in error_msg or "class not registered" in error_msg.lower():
                print("[INFO] PowerPoint not installed, using fallback...")
            else:
                print(f"[INFO] COM error: {error_msg}, using fallback...")
            return ppt_to_images_fallback(ppt_file, output_dir)
    else:
        # Non-Windows: use fallback
        return ppt_to_images_fallback(ppt_file, output_dir)


def list_slides(output_dir="output/slides"):
    """List all slide images"""
    import glob

    slides = sorted(glob.glob(os.path.join(output_dir, "slide*.png")))

    if slides:
        print(f"\n[INFO] Total slides converted: {len(slides)}")
        print(f"[INFO] Location: {output_dir}/")
        for slide in slides:
            size = os.path.getsize(slide) / 1024  # KB
            print(f"   - {os.path.basename(slide)} ({size:.1f} KB)")
    else:
        print("\n[INFO] No slide images found")


if __name__ == "__main__":
    # Default PPT file
    ppt_file = "output/buddha_women_freedom.pptx"

    # Check if custom file is provided
    if len(sys.argv) > 1:
        ppt_file = sys.argv[1]

    # Check if file exists
    if not os.path.exists(ppt_file):
        # Try to find latest PPT in output folder
        import glob
        ppts = sorted(glob.glob("output/*.pptx"), key=os.path.getmtime, reverse=True)

        if ppts:
            ppt_file = ppts[0]
            print(f"[INFO] Using latest PPT: {ppt_file}\n")
        else:
            print("[ERROR] No PowerPoint file found in output folder!")
            print("[INFO] Usage: python ppt_to_images.py [path/to/presentation.pptx]")
            sys.exit(1)

    print("=" * 50)
    print("  PowerPoint to Images Converter")
    print("=" * 50)
    print()

    # Convert
    success = ppt_to_images(ppt_file)

    if success:
        # List generated slides
        list_slides()
        print("\n[OK] Conversion complete!")
    else:
        print("\n[WARNING] Conversion failed")
