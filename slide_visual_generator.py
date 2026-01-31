"""
Automatic Visual Generator for PPT Slides
Adds images, charts, and diagrams based on content type
"""

import os
import re
from typing import Dict, List, Optional, Tuple
from pptx.util import Inches
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import random


class SlideVisualGenerator:
    """
    Analyzes slide content and automatically adds appropriate visuals:
    - Concept/intro slides → Images
    - Data/comparison slides → Charts (bar, pie, line)
    - Timeline/history slides → Timeline diagrams
    """

    def __init__(self):
        self.content_keywords = {
            'chart_bar': ['comparison', 'versus', 'vs', 'difference', 'compare', 'contrast'],
            'chart_pie': ['distribution', 'percentage', 'share', 'proportion', 'breakdown'],
            'chart_line': ['trend', 'growth', 'over time', 'timeline', 'progress', 'evolution'],
            'timeline': ['history', 'timeline', 'chronology', 'events', 'milestones'],
            'table': ['data', 'statistics', 'numbers', 'metrics', 'figures'],
            'image': ['introduction', 'overview', 'concept', 'what is', 'definition']
        }

    def analyze_slide_type(self, title: str, bullets: List[str]) -> str:
        """
        Analyze slide content to determine what type of visual to add

        Returns: 'chart_bar', 'chart_pie', 'chart_line', 'timeline', 'table', 'image', or 'none'
        """
        content_text = f"{title} {' '.join(bullets)}".lower()

        # Check for specific patterns
        for visual_type, keywords in self.content_keywords.items():
            for keyword in keywords:
                if keyword in content_text:
                    return visual_type

        # Default: add image for intro/concept slides
        if any(word in content_text for word in ['introduction', 'overview', 'what', 'about']):
            return 'image'

        # Check if content has numbers (likely needs chart)
        if re.search(r'\d+%|\d+\.\d+', content_text):
            return 'chart_bar'

        return 'none'

    def extract_data_from_bullets(self, bullets: List[str]) -> Tuple[List[str], List[float]]:
        """
        Extract categories and values from bullet points for charts

        Example:
        - "Sales increased by 45%" → ("Sales", 45)
        - "Market share: 23.5%" → ("Market share", 23.5)
        """
        categories = []
        values = []

        for bullet in bullets:
            # Try to extract number and category
            number_match = re.search(r'(\d+\.?\d*)%?', bullet)
            if number_match:
                value = float(number_match.group(1))
                # Get category (text before the number)
                category = bullet[:number_match.start()].strip()
                if not category:
                    category = bullet.split(':')[0] if ':' in bullet else f"Item {len(categories)+1}"

                categories.append(category[:30])  # Limit length
                values.append(value)

        # If no data found, generate sample data
        if not categories:
            categories = [f"Category {i+1}" for i in range(min(4, len(bullets)))]
            values = [random.randint(20, 80) for _ in categories]

        return categories, values

    def add_bar_chart(self, slide, title: str, bullets: List[str]):
        """Add a bar chart to the slide"""
        try:
            categories, values = self.extract_data_from_bullets(bullets)

            # Create chart data
            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series('Values', values)

            # Add chart to slide
            x, y, cx, cy = Inches(5.5), Inches(1.5), Inches(4), Inches(4.5)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart

            chart.has_legend = False
            chart.has_title = True
            chart.chart_title.text_frame.text = title

            return True
        except Exception as e:
            print(f"Error adding bar chart: {e}")
            return False

    def add_pie_chart(self, slide, title: str, bullets: List[str]):
        """Add a pie chart to the slide"""
        try:
            categories, values = self.extract_data_from_bullets(bullets)

            # Normalize to percentages
            total = sum(values)
            if total > 0:
                values = [(v / total) * 100 for v in values]

            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series('Distribution', values)

            x, y, cx, cy = Inches(5.5), Inches(1.5), Inches(4), Inches(4.5)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
            ).chart

            chart.has_legend = True
            chart.chart_title.text_frame.text = title

            return True
        except Exception as e:
            print(f"Error adding pie chart: {e}")
            return False

    def add_line_chart(self, slide, title: str, bullets: List[str]):
        """Add a line chart to the slide"""
        try:
            categories, values = self.extract_data_from_bullets(bullets)

            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series('Trend', values)

            x, y, cx, cy = Inches(5.5), Inches(1.5), Inches(4), Inches(4.5)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
            ).chart

            chart.has_legend = False
            chart.chart_title.text_frame.text = title

            return True
        except Exception as e:
            print(f"Error adding line chart: {e}")
            return False

    def add_visual_to_slide(self, slide, slide_data: dict, use_ai_images: bool = False):
        """
        Main method: Analyzes slide and adds appropriate visual

        Args:
            slide: python-pptx slide object
            slide_data: Dict with 'title' and 'bullets'
            use_ai_images: If True, generate AI images for concept slides
        """
        title = slide_data.get('title', '')
        bullets = slide_data.get('bullets', [])

        # Determine visual type
        visual_type = self.analyze_slide_type(title, bullets)

        print(f"📊 Slide: {title}")
        print(f"   Visual type: {visual_type}")

        # Add appropriate visual
        if visual_type == 'chart_bar':
            return self.add_bar_chart(slide, title, bullets)

        elif visual_type == 'chart_pie':
            return self.add_pie_chart(slide, title, bullets)

        elif visual_type == 'chart_line':
            return self.add_line_chart(slide, title, bullets)

        elif visual_type == 'image' and use_ai_images:
            # Use AI image generation (already implemented in ai_image_generator.py)
            try:
                from ai_image_generator import generate_slide_image
                image_path = generate_slide_image(title, '\n'.join(bullets))
                if image_path:
                    slide.shapes.add_picture(
                        image_path,
                        Inches(5.2), Inches(1.3),
                        width=Inches(4.3), height=Inches(4)
                    )
                    return True
            except Exception as e:
                print(f"   AI image generation failed: {e}")

        elif visual_type == 'timeline':
            # Create simple timeline visual (placeholder for now)
            print(f"   Timeline visual (to be implemented)")

        return False


def enhance_slides_with_visuals(slides: List[dict], use_ai_images: bool = False) -> List[dict]:
    """
    Analyze all slides and add visual metadata

    Args:
        slides: List of slide dictionaries
        use_ai_images: Whether to generate AI images

    Returns:
        Enhanced slides with visual_type metadata
    """
    generator = SlideVisualGenerator()

    enhanced_slides = []
    for idx, slide in enumerate(slides):
        # Skip title slide
        if idx == 0:
            enhanced_slides.append(slide)
            continue

        visual_type = generator.analyze_slide_type(
            slide.get('title', ''),
            slide.get('bullets', [])
        )

        slide['visual_type'] = visual_type
        slide['needs_visual'] = visual_type != 'none'

        enhanced_slides.append(slide)

    return enhanced_slides


# Example usage
if __name__ == "__main__":
    # Test slide analysis
    generator = SlideVisualGenerator()

    test_slides = [
        {
            "title": "Market Share Distribution",
            "bullets": [
                "Company A: 35% market share",
                "Company B: 28% market share",
                "Company C: 22% market share",
                "Others: 15% market share"
            ]
        },
        {
            "title": "Sales Growth Over Time",
            "bullets": [
                "2020: $1.2M revenue",
                "2021: $1.8M revenue",
                "2022: $2.5M revenue",
                "2023: $3.2M revenue"
            ]
        },
        {
            "title": "Introduction to AI",
            "bullets": [
                "Artificial Intelligence is transforming industries",
                "Machine learning enables predictive analytics",
                "Deep learning powers computer vision"
            ]
        }
    ]

    for slide in test_slides:
        visual_type = generator.analyze_slide_type(slide['title'], slide['bullets'])
        print(f"\nSlide: {slide['title']}")
        print(f"Visual Type: {visual_type}")
